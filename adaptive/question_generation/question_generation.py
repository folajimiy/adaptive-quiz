"""
question_generation.py

Hybrid LLM pipeline for generating Java MCQs + adaptive metadata,
restricted to the scope of your lecture slides.

- Uses OpenAI (GPT-4.1 by default) to generate questions + metadata
- Uses Anthropic Claude Sonnet to refine / validate / score items
- Restricts generation using text chunks extracted from .pptx slides
- Writes results to `data/java_questions_adaptive.csv`

Environment variables required:
    OPENAI_API_KEY       - for OpenAI
    ANTHROPIC_API_KEY    - for Anthropic Claude

Directory layout (expected):
    question_generation/
        question_generation.py   <-- this file
        slides/                  <-- put all your .pptx here
        data/
            java_questions_adaptive.csv  <-- created/extended by this script

Run:
    python question_generation.py
"""

import os
import csv
import json
import uuid
import time
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, Tuple

import pandas as pd
from pptx import Presentation

from openai import OpenAI
from anthropic import Anthropic

# ==========================
# CONFIG
# ==========================

OPENAI_MODEL = "gpt-4.1"
CLAUDE_MODEL = "claude-haiku-4-5-20251001"


SLIDES_DIR = os.path.join(os.path.dirname(__file__), "slides")
DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
os.makedirs(DATA_DIR, exist_ok=True)

OUTPUT_CSV = os.path.join(DATA_DIR, "java_questions_adaptive.csv")

DEBUG = True          # minimal transparent debug info
TOP_K_CHUNKS = 5      # number of slide chunks to include as context
MAX_CONTEXT_CHARS = 4000

BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]

TOPIC_MAP = {
    "Java Fundamentals": [
        "Variables and data types",
        "Operators and expressions",
        "Basic syntax and structure",
        "Wrapper classes and autoboxing"
    ],
    "Control Flow": [
        "if/else branching",
        "switch statements",
        "boolean expressions"
    ],
    "Loops": [
        "for loop",
        "while loop",
        "do-while loop",
        "loop control (break/continue)"
    ],
    "Arrays": [
        "1D arrays",
        "Array iteration",
        "Common array errors"
    ],
    "Strings": [
        "String immutability",
        "Common String methods",
        "String comparison and interning"
    ],
    "Methods": [
        "Parameter passing",
        "Method signatures",
        "Return values"
    ],
    "Objects and Classes": [
        "Fields and methods",
        "Constructors",
        "Instance vs static members",
        "toString and equals overrides"
    ],
    "Encapsulation": [
        "Access modifiers",
        "Getters and setters",
        "Information hiding"
    ],
    "Inheritance": [
        "Superclass/subclass relationships",
        "Method overriding",
        "super keyword"
    ],
    "Polymorphism": [
        "Dynamic dispatch",
        "Upcasting and downcasting",
        "Method binding"
    ],
    "Abstract Classes": [
        "Abstract methods",
        "Partial implementation",
        "Concrete subclasses of abstract classes"
    ],
    "Interfaces": [
        "Interface contracts",
        "Implementing multiple interfaces",
        "Functional interfaces"
    ],
    "Generics": [
        "Generic classes",
        "Generic methods",
        "Type parameters and type safety"
    ],
    "Collections": [
        "Lists (e.g., ArrayList)",
        "Sets and uniqueness",
        "Maps (key-value pairs)",
        "Iteration over collections"
    ],
    "JavaFX": [
        "Scene graph",
        "UI controls and layout",
        "Stages and scenes"
    ],
    "Event-Driven Programming": [
        "Event handlers",
        "Listener patterns",
        "JavaFX event handling with lambdas"
    ]
}

# How many questions per (topic, subtopic, bloom) pair
QUESTIONS_PER_PAIR = 3

SUBJECT = "Java"

# ==========================
# LLM CLIENTS
# ==========================

def get_openai_client() -> OpenAI:
    api_key = os.environ.get("OPENAI_API_KEY")
    if not api_key:
        raise RuntimeError("OPENAI_API_KEY environment variable is not set.")
    return OpenAI(api_key=api_key)

def get_claude_client() -> Anthropic:
    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        raise RuntimeError("ANTHROPIC_API_KEY environment variable is not set.")
    return Anthropic(api_key=api_key)

# ==========================
# SLIDE PROCESSING
# ==========================

@dataclass
class SlideChunk:
    slide_file: str
    slide_index: int
    text: str

def extract_text_from_pptx(path: str) -> List[str]:
    prs = Presentation(path)
    texts = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                t = shape.text.strip()
                if t:
                    parts.append(t)
        slide_text = "\n".join(parts).strip()
        if slide_text:
            texts.append(slide_text)
    return texts

def load_slide_chunks(slides_dir: str) -> List[SlideChunk]:
    chunks: List[SlideChunk] = []
    if not os.path.isdir(slides_dir):
        return chunks
    for fname in os.listdir(slides_dir):
        if not fname.lower().endswith(".pptx"):
            continue
        full = os.path.join(slides_dir, fname)
        try:
            slide_texts = extract_text_from_pptx(full)
            for idx, txt in enumerate(slide_texts):
                chunks.append(SlideChunk(slide_file=fname, slide_index=idx, text=txt))
        except Exception as e:
            print(f"⚠️ Failed to read {full}: {e}")
    if DEBUG:
        print(f"Loaded {len(chunks)} slide chunks from {slides_dir}")
    return chunks

def simple_score(query: str, text: str) -> int:
    query_tokens = {t.lower() for t in query.replace(",", " ").split() if len(t) > 2}
    text_tokens = {t.lower() for t in text.replace(",", " ").split() if len(t) > 2}
    return len(query_tokens & text_tokens)

def build_context_for(topic: str,
                      subtopic: str,
                      slide_chunks: List[SlideChunk],
                      top_k: int = TOP_K_CHUNKS,
                      max_chars: int = MAX_CONTEXT_CHARS
                      ) -> Tuple[str, List[Tuple[str, int]]]:
    if not slide_chunks:
        return "", []

    query = f"{topic} - {subtopic}"
    scored = []
    for ch in slide_chunks:
        s = simple_score(query, ch.text)
        if s > 0:
            scored.append((s, ch))
    if not scored:
        # if nothing matched, just take first N chunks
        scored = [(1, ch) for ch in slide_chunks[:top_k]]

    scored.sort(key=lambda x: x[0], reverse=True)
    selected = [ch for _, ch in scored[:top_k]]

    context_parts = []
    debug_chunks = []
    total_chars = 0
    for ch in selected:
        snippet = ch.text.strip()
        if total_chars + len(snippet) > max_chars:
            break
        context_parts.append(snippet)
        total_chars += len(snippet)
        debug_chunks.append((ch.slide_file, ch.slide_index))

    context = "\n\n---\n\n".join(context_parts)
    return context, debug_chunks

# ==========================
# PROMPT BUILDING
# ==========================

def build_generation_prompt(topic: str,
                            subtopic: str,
                            bloom_level: str,
                            subject: str,
                            context: str) -> str:
    return f"""
You are an expert Java instructor and assessment designer.

Your task is to generate ONE high-quality multiple-choice question (MCQ) for an undergraduate {subject} course.

The question MUST:

- Be about the topic: "{topic}"
- Focus specifically on the subtopic: "{subtopic}"
- Operate at Bloom level: "{bloom_level}"
- Be fully contained within the scope of the provided course material (do NOT introduce concepts that clearly require content beyond the slides).
- Be creative, varied, and realistic. Do NOT mention slides, PPT, lecture, or any phrasing like "according to the slides".

Use the following COURSE CONTEXT (extracted from lecture slides) ONLY as background knowledge. You do NOT need to quote it. Just use it as a content boundary:

[COURSE CONTEXT START]
{context}
[COURSE CONTEXT END]

Return your answer as a SINGLE JSON object with the following structure (no extra commentary, no markdown):

{{
  "topic": "{topic}",
  "subtopic": "{subtopic}",
  "subject": "{subject}",
  "bloom_level": "{bloom_level}",
  "question_stem": "...",
  "options": {{
    "A": "...",
    "B": "...",
    "C": "...",
    "D": "..."
  }},
  "correct_answer": "A",
  "a_explanation": "...",
  "b_explanation": "...",
  "c_explanation": "...",
  "d_explanation": "...",
  "main_explanation": "...",

  "cognitive_process": "e.g., recall / explain / apply / analyze / evaluate",
  "kc_tags": ["short", "keyword", "tags"],
  "item_type": "conceptual MCQ or code-based MCQ",

  "predicted_difficulty_level": 3,
  "predicted_difficulty_label": "easy / medium / hard",
  "reasoning_depth": 2,
  "linguistic_complexity": "simple / moderate / complex",
  "estimated_time_seconds": 60,

  "distractor_analysis": "Brief description of how distractors relate to common misunderstandings.",
  "variant_group_id": "short_id_grouping_similar_items",

  "irt_difficulty_b": 0.0,
  "irt_discrimination_a": 1.0,
  "irt_guessing_c": 0.25,

  "remediation_reference": "Natural language hint on what to review.",
  "safety_notes": "",
  "bias_notes": "",

  "misconception_tags_per_option": {{
    "A": ["tag1_if_any"],
    "B": ["tag2_if_any"],
    "C": [],
    "D": []
  }}
}}
"""

def build_claude_refinement_prompt(raw_json: str,
                                   topic: str,
                                   subtopic: str,
                                   bloom_level: str,
                                   subject: str,
                                   context: str) -> str:
    return f"""
You are reviewing a JSON-defined multiple-choice question for an undergraduate {subject} course.

The question was generated to:
- Be about topic: "{topic}"
- Focus on subtopic: "{subtopic}"
- Match Bloom level: "{bloom_level}"
- Stay within the scope of the provided course material

Here is the COURSE CONTEXT extracted from slides (this is the content boundary; do not introduce material clearly beyond this):

[COURSE CONTEXT START]
{context}
[COURSE CONTEXT END]

Here is the ORIGINAL JSON for the question:

[QUESTION JSON START]
{raw_json}
[QUESTION JSON END]

Your tasks:

1. Parse the JSON and check:
   - Technical correctness of the Java content.
   - That the correct option is truly correct and the distractors are plausible but wrong.
   - That the question matches the intended Bloom level.
   - That the content is within the likely scope of the course material.

2. If needed, minimally adjust:
   - question_stem
   - options and correct_answer
   - explanations
   - kc_tags
   - predicted_difficulty_level / label
   - misconception_tags_per_option
   - any other metadata fields

3. Add an evaluation object with:
   - eval_relevance: 1-5
   - eval_bloom_alignment: 1-5
   - eval_accuracy: 1-5
   - eval_explainability: 1-5
   - eval_justification: short justification string

Return a SINGLE JSON object with the same top-level fields as the original question, plus an extra field:

"evaluation": {{
  "eval_relevance": 5,
  "eval_bloom_alignment": 5,
  "eval_accuracy": 5,
  "eval_explainability": 5,
  "eval_justification": "..."
}}

Output ONLY valid JSON. No markdown, no comments, no backticks.
"""

# ==========================
# LLM CALLS
# ==========================

def call_openai_generate(prompt: str, client: OpenAI) -> Dict[str, Any]:
    resp = client.chat.completions.create(
        model=OPENAI_MODEL,
        messages=[
            {"role": "system", "content": "You are a careful Java tutor and assessment designer. Always respond with strict JSON only."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.8,
        max_tokens=800,
    )
    text = resp.choices[0].message.content.strip()
    try:
        return json.loads(text)
    except Exception as e:
        raise ValueError(f"Failed to parse OpenAI JSON: {e}\nRaw text:\n{text}")

def call_claude_refine(raw_item: Dict[str, Any],
                       topic: str,
                       subtopic: str,
                       bloom_level: str,
                       subject: str,
                       context: str,
                       client: Anthropic) -> Dict[str, Any]:
    raw_json = json.dumps(raw_item, ensure_ascii=False, indent=2)
    prompt = build_claude_refinement_prompt(
        raw_json=raw_json,
        topic=topic,
        subtopic=subtopic,
        bloom_level=bloom_level,
        subject=subject,
        context=context
    )
    resp = client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=900,
        temperature=0.2,
        messages=[{"role": "user", "content": prompt}]
    )
    text = resp.content[0].text.strip()
    try:
        return json.loads(text)
    except Exception as e:
        raise ValueError(f"Failed to parse Claude JSON: {e}\nRaw text:\n{text}")

# ==========================
# CSV HANDLING
# ==========================

CSV_COLUMNS = [
    "id",
    "timestamp",
    "topic",
    "subtopic",
    "subject",
    "bloom_level",
    "cognitive_process",
    "kc_tags",
    "question_stem",
    "option_a",
    "option_b",
    "option_c",
    "option_d",
    "correct_answer",
    "a_explanation",
    "b_explanation",
    "c_explanation",
    "d_explanation",
    "main_explanation",
    "item_type",
    "predicted_difficulty_level",
    "predicted_difficulty_label",
    "reasoning_depth",
    "linguistic_complexity",
    "estimated_time_seconds",
    "distractor_analysis",
    "variant_group_id",
    "irt_difficulty_b",
    "irt_discrimination_a",
    "irt_guessing_c",
    "remediation_reference",
    "safety_notes",
    "bias_notes",
    "raw_model_response",
    "retrieved_slide_files",
    "retrieved_chunk_indices",
    "generation_seed_prompt",
    "eval_relevance",
    "eval_bloom_alignment",
    "eval_accuracy",
    "eval_explainability",
    "eval_justification",
    "misconception_tags_per_option"
]

def ensure_output_csv(path: str):
    if not os.path.exists(path):
        os.makedirs(os.path.dirname(path), exist_ok=True)
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(CSV_COLUMNS)

def append_row_to_csv(path: str, row: Dict[str, Any]):
    with open(path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=CSV_COLUMNS)
        writer.writerow(row)

# ==========================
# PIPELINE
# ==========================

def generate_one_item(topic: str,
                      subtopic: str,
                      bloom_level: str,
                      subject: str,
                      slide_chunks: List[SlideChunk],
                      openai_client: OpenAI,
                      claude_client: Anthropic) -> Optional[pd.Series]:
    context, dbg_chunks = build_context_for(topic, subtopic, slide_chunks)

    if DEBUG:
        print(f"\n--- Generating for Topic: {topic} | Subtopic: {subtopic} | Bloom: {bloom_level}")
        print("Slide chunks used:")
        for fname, idx in dbg_chunks:
            print(f"  - {fname} (slide {idx})")

    gen_prompt = build_generation_prompt(
        topic=topic,
        subtopic=subtopic,
        bloom_level=bloom_level,
        subject=subject,
        context=context
    )

    try:
        raw_item = call_openai_generate(gen_prompt, openai_client)
    except Exception as e:
        print(f"❌ OpenAI generation failed: {e}")
        return None

    try:
        refined_item = call_claude_refine(
            raw_item=raw_item,
            topic=topic,
            subtopic=subtopic,
            bloom_level=bloom_level,
            subject=subject,
            context=context,
            client=claude_client
        )
    except Exception as e:
        print(f"❌ Claude refinement failed, falling back to raw item: {e}")
        refined_item = raw_item
        refined_item.setdefault("evaluation", {
            "eval_relevance": None,
            "eval_bloom_alignment": None,
            "eval_accuracy": None,
            "eval_explainability": None,
            "eval_justification": "Claude refinement failed; using raw item."
        })

    evaluation = refined_item.get("evaluation", {})

    options = refined_item.get("options", {})
    mc_tags = refined_item.get("misconception_tags_per_option", {})

    row = {
        "id": str(uuid.uuid4()),
        "timestamp": pd.Timestamp.now().isoformat(),
        "topic": topic,
        "subtopic": subtopic,
        "subject": subject,
        "bloom_level": refined_item.get("bloom_level", bloom_level),
        "cognitive_process": refined_item.get("cognitive_process", ""),
        "kc_tags": "; ".join(refined_item.get("kc_tags", [])),
        "question_stem": refined_item.get("question_stem", ""),
        "option_a": options.get("A", ""),
        "option_b": options.get("B", ""),
        "option_c": options.get("C", ""),
        "option_d": options.get("D", ""),
        "correct_answer": refined_item.get("correct_answer", ""),
        "a_explanation": refined_item.get("a_explanation", ""),
        "b_explanation": refined_item.get("b_explanation", ""),
        "c_explanation": refined_item.get("c_explanation", ""),
        "d_explanation": refined_item.get("d_explanation", ""),
        "main_explanation": refined_item.get("main_explanation", ""),
        "item_type": refined_item.get("item_type", ""),
        "predicted_difficulty_level": refined_item.get("predicted_difficulty_level", None),
        "predicted_difficulty_label": refined_item.get("predicted_difficulty_label", ""),
        "reasoning_depth": refined_item.get("reasoning_depth", None),
        "linguistic_complexity": refined_item.get("linguistic_complexity", ""),
        "estimated_time_seconds": refined_item.get("estimated_time_seconds", None),
        "distractor_analysis": refined_item.get("distractor_analysis", ""),
        "variant_group_id": refined_item.get("variant_group_id", ""),
        "irt_difficulty_b": refined_item.get("irt_difficulty_b", None),
        "irt_discrimination_a": refined_item.get("irt_discrimination_a", None),
        "irt_guessing_c": refined_item.get("irt_guessing_c", None),
        "remediation_reference": refined_item.get("remediation_reference", ""),
        "safety_notes": refined_item.get("safety_notes", ""),
        "bias_notes": refined_item.get("bias_notes", ""),
        "raw_model_response": json.dumps(raw_item, ensure_ascii=False),
        "retrieved_slide_files": "; ".join(sorted({f for f, _ in dbg_chunks})),
        "retrieved_chunk_indices": json.dumps(
            [{"file": f, "slide_index": i} for f, i in dbg_chunks],
            ensure_ascii=False
        ),
        "generation_seed_prompt": gen_prompt,
        "eval_relevance": evaluation.get("eval_relevance"),
        "eval_bloom_alignment": evaluation.get("eval_bloom_alignment"),
        "eval_accuracy": evaluation.get("eval_accuracy"),
        "eval_explainability": evaluation.get("eval_explainability"),
        "eval_justification": evaluation.get("eval_justification", ""),
        "misconception_tags_per_option": json.dumps(mc_tags, ensure_ascii=False)
    }

    return pd.Series(row)

def run_generation():
    print("🔧 Loading slide chunks...")
    slide_chunks = load_slide_chunks(SLIDES_DIR)
    if not slide_chunks:
        print(f"⚠️ No slides found in {SLIDES_DIR}. The model will still run, but without slide restriction.")

    ensure_output_csv(OUTPUT_CSV)

    openai_client = get_openai_client()
    claude_client = get_claude_client()

    total_generated = 0

    for topic, subtopics in TOPIC_MAP.items():
        for subtopic in subtopics:
            for bloom_level in BLOOM_LEVELS:
                for i in range(QUESTIONS_PER_PAIR):
                    print(f"\n=== {topic} | {subtopic} | {bloom_level} | Q#{i+1} ===")
                    try:
                        series = generate_one_item(
                            topic=topic,
                            subtopic=subtopic,
                            bloom_level=bloom_level,
                            subject=SUBJECT,
                            slide_chunks=slide_chunks,
                            openai_client=openai_client,
                            claude_client=claude_client
                        )
                        if series is not None:
                            append_row_to_csv(OUTPUT_CSV, series.to_dict())
                            total_generated += 1
                            print(f"✅ Saved question {total_generated} to CSV.")
                        else:
                            print("❌ Skipped due to generation failure.")
                    except Exception as e:
                        print(f"🔥 Critical error on {topic} / {subtopic} / {bloom_level}: {e}")
                    time.sleep(1.0)  # gentle pacing

    print(f"\n🎉 Generation complete! Total questions generated this run: {total_generated}")
    print(f"Output CSV: {OUTPUT_CSV}")

if __name__ == "__main__":
    run_generation()
