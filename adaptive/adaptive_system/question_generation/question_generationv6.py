#!/usr/bin/env python3
"""
question_generation_v5.py

Research-grade, Bloom-aware, embedding-based Java MCQ generator
for a novel adaptive system.

Key features:
- Uses OpenAI embeddings + your .pptx slides to bound question scope to your course.
- Uses OpenAI GPT-4.1-mini to generate creative MCQs (no slide references).
- Randomizes correct answer position (A/B/C/D) in code.
- Uses Claude 3.5 Sonnet to generate rich psychometric + pedagogical metadata:
  * cognitive_process
  * predicted_difficulty_level / label
  * reasoning_depth
  * linguistic_complexity
  * estimated_time_seconds
  * distractor_analysis
  * variant_group_id
  * IRT parameters (b/a/c)
  * remediation_reference
  * safety_notes
  * bias_notes
  * misconception_tags_per_option per option
  * eval_relevance, eval_bloom_alignment, eval_accuracy, eval_explainability, eval_justification

Outputs a CSV compatible with the adaptive Streamlit tutor
(and rich enough for research / adaptive modeling).
"""

import os
import re
import csv
import json
import uuid
import time
from datetime import datetime
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, Tuple

import numpy as np
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from anthropic import Anthropic
import random

# ==========================
# CONFIG
# ==========================

DEBUG = True  # minimal debug: show which slide files are used

OPENAI_CHAT_MODEL = "gpt-5.1"
OPENAI_EMBED_MODEL = "text-embedding-3-small"
CLAUDE_MODEL = "claude-3-5-haiku-20241022"

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(BASE_DIR, "slides")
DATA_DIR = os.path.join(BASE_DIR, "data")
os.makedirs(DATA_DIR, exist_ok=True)

OUTPUT_CSV = os.path.join(DATA_DIR, "java_questions_adaptive_gpt5_contd_10_per_each.csv")

# Context control
TOP_K_CHUNKS = 5
MAX_CONTEXT_CHARS = 4000

SUBJECT = "Java"
QUESTIONS_PER_PAIR = 10

BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]

TOPIC_MAP = {
    # "Java Fundamentals": [
    #     "Variables and data types",
    #     "Operators and expressions",
    #     "Basic syntax and structure",
    #     "Wrapper classes and autoboxing"
    # ],
    # "Control Flow": [
    #     "if/else branching",
    #     "switch statements",
    #     "boolean expressions"
    # ],
    # "Loops": [
    #     "for loop",
    #     "while loop",
    #     "do-while loop",
    #     "loop control (break/continue)"
    # ],
    # "Arrays": [
    #     "1D arrays",
    #     "Array iteration",
    #     "Common array errors"
    # ],
    # "Strings": [
    #     "String immutability",
    #     "Common String methods",
    #     "String comparison and interning"
    # ],
    # "Methods": [
    #     "Parameter passing",
    #     "Method signatures",
    #     "Return values"
    # ],
    "Objects and Classes": [
        "Fields and methods",
        "Constructors",
        "Instance vs static members",
        "toString and equals overrides"
    ],
    "Encapsulation": [
        "Access modifiers",
        "Getters and setters",
        "Information hiding"
    ],
    "Inheritance": [
        "Superclass/subclass relationships",
        "Method overriding",
        "super keyword"
    ],
    "Polymorphism": [
        "Dynamic dispatch",
        "Upcasting and downcasting",
        "Method binding"
    ],
    "Abstract Classes": [
        "Abstract methods",
        "Partial implementation",
        "Concrete subclasses of abstract classes"
    ],
    "Interfaces": [
        "Interface contracts",
        "Implementing multiple interfaces",
        "Functional interfaces"
    ],
    "Generics": [
        "Generic classes",
        "Generic methods",
        "Type parameters and type safety"
    ],
    "Collections": [
        "Lists (e.g., ArrayList)",
        "Sets and uniqueness",
        "Maps (key-value pairs)",
        "Iteration over collections"
    ],
    "JavaFX": [
        "Scene graph",
        "UI controls and layout",
        "Stages and scenes"
    ],
    "Event-Driven Programming": [
        "Event handlers",
        "Listener patterns",
        "JavaFX event handling with lambdas"
    ]
}

# CSV schema (matches earlier adaptive tutor + adds psychometric metadata)
CSV_COLUMNS = [
    "id",
    "timestamp",
    "topic",
    "subtopic",
    "subject",
    "bloom_level",
    "cognitive_process",
    "kc_tags",
    "question_stem",
    "option_a",
    "option_b",
    "option_c",
    "option_d",
    "correct_answer",
    "a_explanation",
    "b_explanation",
    "c_explanation",
    "d_explanation",
    "main_explanation",
    "item_type",
    "predicted_difficulty_level",
    "predicted_difficulty_label",
    "reasoning_depth",
    "linguistic_complexity",
    "estimated_time_seconds",
    "distractor_analysis",
    "variant_group_id",
    "irt_difficulty_b",
    "irt_discrimination_a",
    "irt_guessing_c",
    "remediation_reference",
    "safety_notes",
    "bias_notes",
    "raw_model_response",
    "retrieved_slide_files",
    "retrieved_chunk_indices",
    "generation_seed_prompt",
    "eval_relevance",
    "eval_bloom_alignment",
    "eval_accuracy",
    "eval_explainability",
    "eval_justification",
    "misconception_tags_per_option"
]

# ==========================
# CLIENTS
# ==========================

def init_clients() -> Tuple[OpenAI, Anthropic]:
    openai_key = os.getenv("OPENAI_API_KEY")
    claude_key = os.getenv("ANTHROPIC_API_KEY") or os.getenv("CLAUDE_ANTHROPIC_KEY")

    if not openai_key:
        raise RuntimeError("OPENAI_API_KEY not set.")
    if not claude_key:
        raise RuntimeError("ANTHROPIC_API_KEY or CLAUDE_ANTHROPIC_KEY not set.")

    openai_client = OpenAI(api_key=openai_key)
    anthropic_client = Anthropic(api_key=claude_key)
    return openai_client, anthropic_client

# ==========================
# SLIDES + EMBEDDINGS
# ==========================

@dataclass
class SlideChunk:
    text: str
    slide_file: str
    slide_index: int

def extract_text_from_pptx(path: str) -> List[str]:
    prs = Presentation(path)
    slide_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                t = (shape.text or "").strip()
                if t:
                    parts.append(t)
        merged = "\n".join(parts).strip()
        if merged:
            slide_texts.append(merged)
    return slide_texts

def load_slide_chunks(slides_dir: str) -> List[SlideChunk]:
    chunks: List[SlideChunk] = []
    if not os.path.isdir(slides_dir):
        print(f"⚠️ Slides directory not found: {slides_dir}")
        return chunks

    for fname in os.listdir(slides_dir):
        if not fname.lower().endswith(".pptx"):
            continue
        full = os.path.join(slides_dir, fname)
        try:
            texts = extract_text_from_pptx(full)
            for idx, txt in enumerate(texts):
                chunks.append(SlideChunk(txt, fname, idx))
        except Exception as e:
            print(f"⚠️ Failed to read {full}: {e}")

    print(f"📌 Loaded {len(chunks)} slide chunks from {slides_dir}")
    return chunks

def embed_texts(openai_client: OpenAI, texts: List[str]) -> List[np.ndarray]:
    if not texts:
        return []
    resp = openai_client.embeddings.create(
        model=OPENAI_EMBED_MODEL,
        input=texts
    )
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]

def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    denom = float(np.linalg.norm(a) * np.linalg.norm(b))
    if denom == 0.0:
        return 0.0
    return float(np.dot(a, b) / denom)

def get_relevant_slide_context(
    openai_client: OpenAI,
    slide_chunks: List[SlideChunk],
    slide_embeddings: List[np.ndarray],
    query: str,
    k: int = TOP_K_CHUNKS,
    max_chars: int = MAX_CONTEXT_CHARS
) -> Tuple[str, List[Tuple[str, int]]]:
    """
    Returns:
      context_text: concatenated slide snippets (used as content boundary)
      debug_info: list of (file, slide_index) used
    """
    if not slide_chunks or not slide_embeddings:
        return "", []

    query_emb = embed_texts(openai_client, [query])[0]
    sims = [cosine_similarity(query_emb, emb) for emb in slide_embeddings]
    top_idx = np.argsort(sims)[-k:][::-1]

    chunks = []
    debug_info: List[Tuple[str, int]] = []
    total_chars = 0

    for i in top_idx:
        ch = slide_chunks[i]
        txt = ch.text.strip()
        if not txt:
            continue
        if total_chars + len(txt) > max_chars:
            break
        chunks.append(txt)
        total_chars += len(txt)
        debug_info.append((ch.slide_file, ch.slide_index))

    context = "\n\n---\n\n".join(chunks)
    return context, debug_info

# ==========================
# PROMPTS
# ==========================

def build_generation_prompt(
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str,
    context: str
) -> str:
    """
    Generation prompt: slide-bounded but creative and Bloom-aware.
    """
    return f"""
You are generating an ORIGINAL multiple-choice question for a university-level {subject} course.

Constraints:
- Topic: "{topic}"
- Subtopic: "{subtopic}"
- Bloom level: "{bloom_level}"
- The question MUST be consistent with an intermediate object-oriented Java course.
- Use the course context below ONLY as a content boundary.
  * Do NOT quote it.
  * Do NOT say 'according to the slides' or mention slides/lectures.
  * Just ensure you do not go beyond the ideas/framing suggested by it.

[COURSE CONTEXT - DO NOT QUOTE DIRECTLY]
{context}
[END CONTEXT]

Question requirements:
1. Stem must be clear, focused, and ≤ 50 words.
2. Exactly 4 options: A, B, C, D (one correct, three distractors).
3. Crafted to match the cognitive demand of the given Bloom level.
4. Distractors should reflect realistic student misconceptions.

Format your response EXACTLY like this (no extra commentary):

Question: [Your question stem here]

Options:
A) [Option A text]
B) [Option B text]
C) [Option C text]
D) [Option D text]

Correct Answer: [Letter of correct option only - A, B, C, or D]

Option A explanation: [Brief explanation of why option A is correct or incorrect]
Option B explanation: [Brief explanation of why option B is correct or incorrect]
Option C explanation: [Brief explanation of why option C is correct or incorrect]
Option D explanation: [Brief explanation of why option D is correct or incorrect]

Main Explanation: [Brief overall explanation of the correct answer]

Tags: [comma-separated list of short conceptual tags]
"""

def build_claude_metadata_prompt(
    item: Dict[str, Any],
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str
) -> str:
    """
    Ask Claude to generate psychometric + pedagogical metadata in JSON.
    """
    return f"""
You are analyzing a Java multiple-choice question for an adaptive learning system.

Here is the item (AFTER answer choice randomization):

Question: \"\"\"{item['question']}\"\"\"
Options:
A) {item['option_a']}
B) {item['option_b']}
C) {item['option_c']}
D) {item['option_d']}

Correct Answer: {item['correct_answer']}

Option A explanation: {item.get('option_a_explanation', '')}
Option B explanation: {item.get('option_b_explanation', '')}
Option C explanation: {item.get('option_c_explanation', '')}
Option D explanation: {item.get('option_d_explanation', '')}
Main Explanation: {item.get('main_explanation', '')}

Metadata context:
- Subject: {subject}
- Topic: {topic}
- Subtopic: {subtopic}
- Intended Bloom level: {bloom_level}

Your tasks:
1. Infer the cognitive_process (e.g., recall, explain, interpret, apply, analyze, evaluate).
2. Refine / extend the KC tags (knowledge components) for this item.
3. Characterize the item psychometrically and pedagogically:
   - predicted_difficulty_level: integer 1 (very easy) to 5 (very hard)
   - predicted_difficulty_label: "very easy", "easy", "medium", "hard", "very hard"
   - reasoning_depth: integer 1-5
   - linguistic_complexity: "simple", "moderate", or "complex"
   - estimated_time_seconds: typical time to answer (e.g., 30-120)
   - item_type: e.g., "conceptual MCQ", "code-based MCQ", "debugging MCQ", "design MCQ"
   - distractor_analysis: brief natural language description of how distractors relate to misconceptions
   - variant_group_id: a short stable identifier for this conceptual family (e.g., "arrays_indexing_off_by_one")

4. Provide approximate IRT-style parameters (heuristic estimates are fine):
   - irt_difficulty_b: real number (negative easier, positive harder)
   - irt_discrimination_a: positive real (e.g., 0.25 to 2.5)
   - irt_guessing_c: between 0.0 and 0.25 for 4-option MCQs

5. Provide remediation and safety:
   - remediation_reference: short text describing what to review (e.g., "Review how equals() differs from == for Strings.")
   - safety_notes: mention anything potentially problematic (usually "" for normal Java CS questions)
   - bias_notes: mention any fairness/issues (e.g., culturally specific names) or "" if none.

6. Provide misconception tags for each option:
   - misconception_tags_per_option: a mapping from option letter to a LIST of short, snake_case tags
     Example:
       "misconception_tags_per_option": {{
         "A": ["string_is_primitive"],
         "B": ["confuse_assignment_with_equality"],
         "C": [],
         "D": []
       }}
     Use tags that will help an instructor reason about common student mistakes.

7. Evaluate the item along 4 dimensions, each 1 (very poor) to 5 (excellent):
   - eval_relevance
   - eval_bloom_alignment
   - eval_accuracy
   - eval_explainability
   Also provide a short eval_justification string.

Return ONLY a single JSON object with the following fields:

{{
  "cognitive_process": "string",
  "kc_tags": ["tag1", "tag2"],
  "item_type": "string",
  "predicted_difficulty_level": 3,
  "predicted_difficulty_label": "medium",
  "reasoning_depth": 2,
  "linguistic_complexity": "moderate",
  "estimated_time_seconds": 60,
  "distractor_analysis": "string",
  "variant_group_id": "string",
  "irt_difficulty_b": 0.0,
  "irt_discrimination_a": 1.0,
  "irt_guessing_c": 0.25,
  "remediation_reference": "string",
  "safety_notes": "string",
  "bias_notes": "string",
  "misconception_tags_per_option": {{
    "A": ["tag_a1", "tag_a2"],
    "B": ["tag_b1"],
    "C": [],
    "D": []
  }},
  "eval_relevance": 5,
  "eval_bloom_alignment": 5,
  "eval_accuracy": 5,
  "eval_explainability": 5,
  "eval_justification": "short justification"
}}

No markdown, no comments, no extra text. JSON ONLY.
"""

# ==========================
# GENERATION (OPENAI)
# ==========================

def generate_with_openai(openai_client: OpenAI, prompt: str) -> str:
    resp = openai_client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a rigorous Java exam question generator. "
                    "You must strictly follow the response format and rules."
                )
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.9,
        top_p=0.95,
        max_completion_tokens=600,
    )
    return resp.choices[0].message.content.strip()

# ==========================
# PARSE MCQ TEXT
# ==========================

def parse_question_response(response: str) -> Dict[str, Any]:
    patterns = {
        'question': r"[Qq]uestion:\s*(.+?)\s*[Oo]ptions?:",
        'option_a': r"A\)\s*(.+?)(?=\n[B-D]\)|\nCorrect Answer:|$)",
        'option_b': r"B\)\s*(.+?)(?=\n[C-D]\)|\nCorrect Answer:|$)",
        'option_c': r"C\)\s*(.+?)(?=\nD\)|\nCorrect Answer:|$)",
        'option_d': r"D\)\s*(.+?)(?=\nCorrect Answer:|$)",
        'correct_answer': r"Correct Answer:\s*([A-D])",
        'option_a_explanation': r"Option A explanation:\s*(.+?)(?=\nOption B explanation:|\nMain Explanation:|$)",
        'option_b_explanation': r"Option B explanation:\s*(.+?)(?=\nOption C explanation:|\nMain Explanation:|$)",
        'option_c_explanation': r"Option C explanation:\s*(.+?)(?=\nOption D explanation:|\nMain Explanation:|$)",
        'option_d_explanation': r"Option D explanation:\s*(.+?)(?=\nMain Explanation:|$)",
        'main_explanation': r"Main Explanation:\s*(.+?)(?=\nTags:|$)",
        'tags': r"Tags:\s*(.+)",
    }

    extracted: Dict[str, Any] = {}
    for key, pattern in patterns.items():
        match = re.search(pattern, response, re.DOTALL)
        extracted[key] = match.group(1).strip() if match else ""

    tags = [t.strip() for t in extracted.get("tags", "").split(",") if t.strip()]
    extracted["tags_list"] = tags
    return extracted

# ==========================
# SHUFFLE OPTIONS
# ==========================

def shuffle_question_structure(q: Dict[str, Any]) -> Dict[str, Any]:
    """
    Randomize A/B/C/D while preserving correctness + explanations.
    """
    options = {
        "A": q.get("option_a", ""),
        "B": q.get("option_b", ""),
        "C": q.get("option_c", ""),
        "D": q.get("option_d", ""),
    }
    expl = {
        "A": q.get("option_a_explanation", ""),
        "B": q.get("option_b_explanation", ""),
        "C": q.get("option_c_explanation", ""),
        "D": q.get("option_d_explanation", ""),
    }
    correct = q.get("correct_answer", "A").strip().upper() or "A"

    original_labels = ["A", "B", "C", "D"]
    random.shuffle(original_labels)

    new_options = {}
    new_expl = {}
    new_correct = None
    new_labels = ["A", "B", "C", "D"]

    for new_label, old_label in zip(new_labels, original_labels):
        new_options[new_label] = options.get(old_label, "")
        new_expl[new_label] = expl.get(old_label, "")
        if old_label == correct:
            new_correct = new_label

    q["option_a"] = new_options["A"]
    q["option_b"] = new_options["B"]
    q["option_c"] = new_options["C"]
    q["option_d"] = new_options["D"]

    q["option_a_explanation"] = new_expl["A"]
    q["option_b_explanation"] = new_expl["B"]
    q["option_c_explanation"] = new_expl["C"]
    q["option_d_explanation"] = new_expl["D"]

    q["correct_answer"] = new_correct or "A"
    return q

# ==========================
# CLAUDE METADATA
# ==========================

def get_claude_metadata(
    anthropic_client: Anthropic,
    item: Dict[str, Any],
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str
) -> Optional[Dict[str, Any]]:
    prompt = build_claude_metadata_prompt(item, topic, subtopic, bloom_level, subject)
    try:
        resp = anthropic_client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=900,
            temperature=0.2,
            messages=[{"role": "user", "content": prompt}],
        )
        text = resp.content[0].text.strip()
        meta = json.loads(text)
        return meta
    except Exception as e:
        print("⚠️ Claude metadata JSON parse/HTTP error:", e)
        return None

# ==========================
# CSV HELPERS
# ==========================

def ensure_output_csv(path: str):
    if not os.path.exists(path):
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(CSV_COLUMNS)

def append_row_to_csv(path: str, row: Dict[str, Any]):
    with open(path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=CSV_COLUMNS)
        writer.writerow(row)

# ==========================
# PIPELINE FOR ONE ITEM
# ==========================

def run_pipeline_for_pair(
    openai_client: OpenAI,
    anthropic_client: Anthropic,
    slide_chunks: List[SlideChunk],
    slide_embeddings: List[np.ndarray],
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str,
    csv_path: str
) -> Optional[Dict[str, Any]]:
    # 1) Retrieve slide context
    query = f"{topic} - {subtopic} - {subject}"
    context, dbg_chunks = get_relevant_slide_context(
        openai_client,
        slide_chunks,
        slide_embeddings,
        query,
        k=TOP_K_CHUNKS,
        max_chars=MAX_CONTEXT_CHARS
    )

    if DEBUG:
        print(f"\nDEBUG: Topic '{topic}' | Subtopic '{subtopic}' | Bloom '{bloom_level}'")
        if dbg_chunks:
            print("DEBUG: Slide chunks used for scope:")
            for f, idx in dbg_chunks:
                print(f"  - {f} (slide {idx})")
        else:
            print("DEBUG: No slide chunks strongly matched; using unspecific generation.")
        print()

    # 2) Build generation prompt
    gen_prompt = build_generation_prompt(
        topic=topic,
        subtopic=subtopic,
        bloom_level=bloom_level,
        subject=subject,
        context=context
    )

    # 3) Generate raw question with OpenAI
    raw_text = generate_with_openai(openai_client, gen_prompt)
    q = parse_question_response(raw_text)

    required = ["question", "option_a", "option_b", "option_c", "option_d", "correct_answer"]
    if not all(q.get(k, "").strip() for k in required):
        print("❌ Incomplete question structure, skipping.")
        return None

    # 4) Shuffle options
    q = shuffle_question_structure(q)

    # Build item dict for metadata model
    item_for_meta = {
        "question": q["question"],
        "option_a": q["option_a"],
        "option_b": q["option_b"],
        "option_c": q["option_c"],
        "option_d": q["option_d"],
        "correct_answer": q["correct_answer"],
        "option_a_explanation": q.get("option_a_explanation", ""),
        "option_b_explanation": q.get("option_b_explanation", ""),
        "option_c_explanation": q.get("option_c_explanation", ""),
        "option_d_explanation": q.get("option_d_explanation", ""),
        "main_explanation": q.get("main_explanation", "")
    }

    # 5) Get metadata from Claude
    meta = get_claude_metadata(
        anthropic_client=anthropic_client,
        item=item_for_meta,
        topic=topic,
        subtopic=subtopic,
        bloom_level=bloom_level,
        subject=subject
    )

    # Fallback heuristic metadata if Claude fails
    if not meta:
        meta = {
            "cognitive_process": bloom_level.lower(),
            "kc_tags": q.get("tags_list", []),
            "item_type": "conceptual MCQ",
            "predicted_difficulty_level": 3,
            "predicted_difficulty_label": "medium",
            "reasoning_depth": 2,
            "linguistic_complexity": "moderate",
            "estimated_time_seconds": 60,
            "distractor_analysis": "",
            "variant_group_id": f"{topic.lower().replace(' ','_')}_{subtopic.lower().replace(' ','_')}",
            "irt_difficulty_b": 0.0,
            "irt_discrimination_a": 1.0,
            "irt_guessing_c": 0.25,
            "remediation_reference": "",
            "safety_notes": "",
            "bias_notes": "",
            "misconception_tags_per_option": {
                "A": [],
                "B": [],
                "C": [],
                "D": []
            },
            "eval_relevance": None,
            "eval_bloom_alignment": None,
            "eval_accuracy": None,
            "eval_explainability": None,
            "eval_justification": "Fallback heuristic metadata; Claude failed."
        }
    else:
        # ensure kc_tags includes question tags from generation as well
        base_tags = q.get("tags_list", [])
        extra_tags = meta.get("kc_tags", [])
        # merge while preserving uniqueness
        merged_tags = list(dict.fromkeys(base_tags + extra_tags))
        meta["kc_tags"] = merged_tags

    # 6) Build row for CSV
    chunk_meta = [{"file": f, "slide_index": i} for f, i in dbg_chunks]
    retrieved_files = sorted({f for f, _ in dbg_chunks})
    retrieved_chunk_indices = json.dumps(chunk_meta, ensure_ascii=False)

    row = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.now().isoformat(),
        "topic": topic,
        "subtopic": subtopic,
        "subject": subject,
        "bloom_level": bloom_level,
        "cognitive_process": meta.get("cognitive_process", ""),
        "kc_tags": "; ".join(meta.get("kc_tags", [])),
        "question_stem": q["question"],
        "option_a": q["option_a"],
        "option_b": q["option_b"],
        "option_c": q["option_c"],
        "option_d": q["option_d"],
        "correct_answer": q["correct_answer"],
        "a_explanation": q.get("option_a_explanation", ""),
        "b_explanation": q.get("option_b_explanation", ""),
        "c_explanation": q.get("option_c_explanation", ""),
        "d_explanation": q.get("option_d_explanation", ""),
        "main_explanation": q.get("main_explanation", ""),
        "item_type": meta.get("item_type", ""),
        "predicted_difficulty_level": meta.get("predicted_difficulty_level", None),
        "predicted_difficulty_label": meta.get("predicted_difficulty_label", ""),
        "reasoning_depth": meta.get("reasoning_depth", None),
        "linguistic_complexity": meta.get("linguistic_complexity", ""),
        "estimated_time_seconds": meta.get("estimated_time_seconds", None),
        "distractor_analysis": meta.get("distractor_analysis", ""),
        "variant_group_id": meta.get("variant_group_id", ""),
        "irt_difficulty_b": meta.get("irt_difficulty_b", None),
        "irt_discrimination_a": meta.get("irt_discrimination_a", None),
        "irt_guessing_c": meta.get("irt_guessing_c", None),
        "remediation_reference": meta.get("remediation_reference", ""),
        "safety_notes": meta.get("safety_notes", ""),
        "bias_notes": meta.get("bias_notes", ""),
        "raw_model_response": raw_text,
        "retrieved_slide_files": "; ".join(retrieved_files),
        "retrieved_chunk_indices": retrieved_chunk_indices,
        "generation_seed_prompt": gen_prompt,
        "eval_relevance": meta.get("eval_relevance", None),
        "eval_bloom_alignment": meta.get("eval_bloom_alignment", None),
        "eval_accuracy": meta.get("eval_accuracy", None),
        "eval_explainability": meta.get("eval_explainability", None),
        "eval_justification": meta.get("eval_justification", ""),
        "misconception_tags_per_option": json.dumps(
            meta.get("misconception_tags_per_option", {}),
            ensure_ascii=False
        )
    }

    append_row_to_csv(csv_path, row)
    return row

# ==========================
# MAIN
# ==========================

def main():
    openai_client, anthropic_client = init_clients()

    slide_chunks = load_slide_chunks(SLIDES_DIR)
    slide_texts = [c.text for c in slide_chunks]
    slide_embeddings = embed_texts(openai_client, slide_texts) if slide_texts else []

    ensure_output_csv(OUTPUT_CSV)

    total = 0

    for topic, subtopics in TOPIC_MAP.items():
        for subtopic in subtopics:
            for bloom in BLOOM_LEVELS:
                for i in range(QUESTIONS_PER_PAIR):
                    print(f"\n>> {topic} | {subtopic} | {bloom} | Q{i+1}")
                    try:
                        row = run_pipeline_for_pair(
                            openai_client=openai_client,
                            anthropic_client=anthropic_client,
                            slide_chunks=slide_chunks,
                            slide_embeddings=slide_embeddings,
                            topic=topic,
                            subtopic=subtopic,
                            bloom_level=bloom,
                            subject=SUBJECT,
                            csv_path=OUTPUT_CSV
                        )
                        if row:
                            total += 1
                            print("✅ Question generated and saved.")
                        else:
                            print("❌ Question generation failed for this combo.")
                    except Exception as e:
                        print("🔥 Critical error:", e)
                    time.sleep(1.5)  # be nice to APIs

    print(f"\n🎉 Done! Total questions generated this run: {total}")
    print("📄 Output CSV:", OUTPUT_CSV)


if __name__ == "__main__":
    main()