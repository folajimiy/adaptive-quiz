
#!/usr/bin/env python3
"""
question_generation.py

Embedding-based, slide-bounded, creative Java MCQ generator with topics + subtopics.

- Automatically loads all .pptx slides from a ./slides folder
- Uses OpenAI embeddings to scope (topic + subtopic) to relevant slides
- Does NOT show slide text to the LLM (slide-bounded, not slide-anchored)
- Generates creative MCQs using OpenAI (GPT-4.1 family)
- Evaluates them using Claude 3.5 Sonnet
- Randomizes correct answer position (A/B/C/D) in code
- Logs topic, subtopic, Bloom level, slide filenames and eval scores into a CSV

Requirements:
    pip install python-pptx openai anthropic numpy pandas

Environment:
    OPENAI_API_KEY         - OpenAI key
    ANTHROPIC_API_KEY      - Anthropic key (or CLAUDE_ANTHROPIC_KEY fallback)
"""

import os
import re
import csv
import json
import uuid
import time
from datetime import datetime
from dataclasses import dataclass
from typing import List, Dict, Any, Optional, Tuple

import numpy as np
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from anthropic import Anthropic
import random

# ==========================
# CONFIG
# ==========================

# Minimal debug: show which slide files are used per (topic, subtopic)
DEBUG = True

# OpenAI + Anthropic models
OPENAI_CHAT_MODEL = "gpt-4.1-mini"
OPENAI_EMBED_MODEL = "text-embedding-3-small"
CLAUDE_MODEL = "claude-haiku-4-5-20251001"

# Paths
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDES_DIR = os.path.join(BASE_DIR, "slides")
DATA_DIR = os.path.join(BASE_DIR, "data")
os.makedirs(DATA_DIR, exist_ok=True)

OUTPUT_CSV = os.path.join(DATA_DIR, "java_questions_adaptive.csv")

# Generation parameters
QUESTIONS_PER_PAIR = 3  # how many questions per (topic, subtopic, bloom)
SUBJECT = "Java"

BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]

TOPIC_MAP = {
    "Java Fundamentals": [
        "Variables and data types",
        "Operators and expressions",
        "Basic syntax and structure",
        "Wrapper classes and autoboxing"
    ],
    "Control Flow": [
        "if/else branching",
        "switch statements",
        "boolean expressions"
    ],
    "Loops": [
        "for loop",
        "while loop",
        "do-while loop",
        "loop control (break/continue)"
    ],
    "Arrays": [
        "1D arrays",
        "Array iteration",
        "Common array errors"
    ],
    "Strings": [
        "String immutability",
        "Common String methods",
        "String comparison and interning"
    ],
    "Methods": [
        "Parameter passing",
        "Method signatures",
        "Return values"
    ],
    "Objects and Classes": [
        "Fields and methods",
        "Constructors",
        "Instance vs static members",
        "toString and equals overrides"
    ],
    "Encapsulation": [
        "Access modifiers",
        "Getters and setters",
        "Information hiding"
    ],
    "Inheritance": [
        "Superclass/subclass relationships",
        "Method overriding",
        "super keyword"
    ],
    "Polymorphism": [
        "Dynamic dispatch",
        "Upcasting and downcasting",
        "Method binding"
    ],
    "Abstract Classes": [
        "Abstract methods",
        "Partial implementation",
        "Concrete subclasses of abstract classes"
    ],
    "Interfaces": [
        "Interface contracts",
        "Implementing multiple interfaces",
        "Functional interfaces"
    ],
    "Generics": [
        "Generic classes",
        "Generic methods",
        "Type parameters and type safety"
    ],
    "Collections": [
        "Lists (e.g., ArrayList)",
        "Sets and uniqueness",
        "Maps (key-value pairs)",
        "Iteration over collections"
    ],
    "JavaFX": [
        "Scene graph",
        "UI controls and layout",
        "Stages and scenes"
    ],
    "Event-Driven Programming": [
        "Event handlers",
        "Listener patterns",
        "JavaFX event handling with lambdas"
    ]
}

# CSV schema (compatible with Streamlit tutor expectations + extra eval fields)
CSV_COLUMNS = [
    "id",
    "timestamp",
    "topic",
    "subtopic",
    "subject",
    "bloom_level",
    "cognitive_process",
    "kc_tags",
    "question_stem",
    "option_a",
    "option_b",
    "option_c",
    "option_d",
    "correct_answer",
    "a_explanation",
    "b_explanation",
    "c_explanation",
    "d_explanation",
    "main_explanation",
    "item_type",
    "eval_relevance",
    "eval_bloom_alignment",
    "eval_accuracy",
    "eval_explainability",
    "eval_justification",
    "raw_model_response",
    "retrieved_slide_files",
    "retrieved_chunk_indices",
    "generation_seed_prompt"
]

# ==========================
# API CLIENTS
# ==========================

def init_clients() -> Tuple[OpenAI, Anthropic]:
    openai_key = os.getenv("OPENAI_API_KEY")
    claude_key = os.getenv("ANTHROPIC_API_KEY") or os.getenv("CLAUDE_ANTHROPIC_KEY")

    if not openai_key:
        raise RuntimeError("OPENAI_API_KEY not set in environment.")
    if not claude_key:
        raise RuntimeError("ANTHROPIC_API_KEY or CLAUDE_ANTHROPIC_KEY not set in environment.")

    openai_client = OpenAI(api_key=openai_key)
    anthropic_client = Anthropic(api_key=claude_key)
    return openai_client, anthropic_client

# ==========================
# SLIDE LOADING + EMBEDDINGS
# ==========================

@dataclass
class SlideChunk:
    text: str
    slide_file: str
    slide_index: int

def extract_text_from_pptx(path: str) -> List[str]:
    prs = Presentation(path)
    slide_texts = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.has_text_frame:
                t = (shape.text or "").strip()
                if t:
                    parts.append(t)
        merged = "\n".join(parts).strip()
        if merged:
            slide_texts.append(merged)
    return slide_texts

def load_slide_chunks() -> Tuple[List[SlideChunk], List[str]]:
    chunks: List[SlideChunk] = []
    sources: List[str] = []

    if not os.path.isdir(SLIDES_DIR):
        print(f"⚠️ Slides directory not found: {SLIDES_DIR}")
        return chunks, sources

    for fname in os.listdir(SLIDES_DIR):
        if not fname.lower().endswith(".pptx"):
            continue
        full_path = os.path.join(SLIDES_DIR, fname)
        try:
            slide_texts = extract_text_from_pptx(full_path)
            for idx, txt in enumerate(slide_texts):
                chunks.append(SlideChunk(text=txt, slide_file=fname, slide_index=idx))
                sources.append(fname)
        except Exception as e:
            print(f"⚠️ Failed to read {full_path}: {e}")

    print(f"📌 Loaded {len(chunks)} slide chunks from {SLIDES_DIR}")
    return chunks, sources

def embed_texts(openai_client: OpenAI, texts: List[str]) -> List[np.ndarray]:
    if not texts:
        return []
    resp = openai_client.embeddings.create(
        model=OPENAI_EMBED_MODEL,
        input=texts
    )
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]

def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    denom = float(np.linalg.norm(a) * np.linalg.norm(b))
    if denom == 0.0:
        return 0.0
    return float(np.dot(a, b) / denom)

def get_relevant_slide_scope(
    openai_client: OpenAI,
    slide_chunks: List[SlideChunk],
    slide_embeddings: List[np.ndarray],
    query: str,
    k: int = 4
) -> Tuple[List[str], List[int]]:
    """
    Uses embeddings to see which slide chunks are most related to the query.
    Returns:
        unique_files: sorted list of slide filenames involved
        indices: list of chunk indices chosen
    """
    if not slide_chunks or not slide_embeddings:
        return [], []

    query_emb = embed_texts(openai_client, [query])[0]
    sims = [cosine_similarity(query_emb, emb) for emb in slide_embeddings]
    top_indices = np.argsort(sims)[-k:][::-1]

    selected_files = [slide_chunks[i].slide_file for i in top_indices]
    unique_files = sorted(set(selected_files))

    return unique_files, list(top_indices)

# ==========================
# PROMPT BUILDING (NO SLIDE TEXT)
# ==========================

def build_prompt(topic: str, subtopic: str, bloom_level: str, subject: str) -> str:
    """
    Creative MCQ generation bounded by topic + subtopic and Bloom level.
    Slides are NOT mentioned and slide text is NOT included.
    We will randomize options in code; the model can always put the correct
    answer as A internally.
    """
    return f"""
You are generating a high-quality, original multiple-choice question
for a university-level {subject} course focused on object-oriented programming in Java.

You MUST follow these rules:

1. Stay strictly within the conceptual boundaries of:
   - Topic: "{topic}"
   - Subtopic: "{subtopic}"

2. Do NOT mention slides, lectures, or course materials.
3. Do NOT introduce advanced Java topics beyond standard introductory/
   intermediate OOP material.
4. The question must be creative, non-trivial, and the stem MUST be ≤ 50 words.
5. Target Bloom level: "{bloom_level}".
6. The question must have exactly 4 options: A, B, C, D.
7. Exactly ONE option must be correct.

Format your response EXACTLY like this (no extra commentary):

Question: [Your question stem here]

Options:
A) [Option A text]
B) [Option B text]
C) [Option C text]
D) [Option D text]

Correct Answer: [Letter of correct option only - A, B, or C, or D]

Option A explanation: [Brief explanation of why option A is correct or incorrect]
Option B explanation: [Brief explanation of why option B is correct or incorrect]
Option C explanation: [Brief explanation of why option C is correct or incorrect]
Option D explanation: [Brief explanation of why option D is correct or incorrect]

Main Explanation: [Brief overall explanation of the correct answer]

Tags: [comma-separated list of short conceptual tags]
"""

# ==========================
# OPENAI GENERATION
# ==========================

def generate_with_openai(openai_client: OpenAI, prompt: str) -> str:
    resp = openai_client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a Java exam question generator. "
                    "You must strictly follow the response format and rules."
                )
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.9,
        top_p=0.95,
        max_tokens=500,
    )
    return resp.choices[0].message.content.strip()

# ==========================
# PARSING THE MCQ
# ==========================

def parse_question_response(response: str) -> Dict[str, Any]:
    patterns = {
        'question': r"[Qq]uestion:\s*(.+?)\s*[Oo]ptions?:",
        'option_a': r"A\)\s*(.+?)(?=\n[B-D]\)|\nCorrect Answer:|$)",
        'option_b': r"B\)\s*(.+?)(?=\n[C-D]\)|\nCorrect Answer:|$)",
        'option_c': r"C\)\s*(.+?)(?=\nD\)|\nCorrect Answer:|$)",
        'option_d': r"D\)\s*(.+?)(?=\nCorrect Answer:|$)",
        'correct_answer': r"Correct Answer:\s*([A-D])",
        'option_a_explanation': r"Option A explanation:\s*(.+?)(?=\nOption B explanation:|\nMain Explanation:|$)",
        'option_b_explanation': r"Option B explanation:\s*(.+?)(?=\nOption C explanation:|\nMain Explanation:|$)",
        'option_c_explanation': r"Option C explanation:\s*(.+?)(?=\nOption D explanation:|\nMain Explanation:|$)",
        'option_d_explanation': r"Option D explanation:\s*(.+?)(?=\nMain Explanation:|$)",
        'main_explanation': r"Main Explanation:\s*(.+?)(?=\nTags:|$)",
        'tags': r"Tags:\s*(.+)",
    }

    extracted: Dict[str, Any] = {}
    for key, pattern in patterns.items():
        match = re.search(pattern, response, re.DOTALL)
        extracted[key] = match.group(1).strip() if match else ""

    # Normalize tags into list
    tags = [t.strip() for t in extracted.get("tags", "").split(",") if t.strip()]
    extracted["tags_list"] = tags

    return extracted

# ==========================
# OPTION SHUFFLING
# ==========================

def shuffle_question_structure(q: Dict[str, Any]) -> Dict[str, Any]:
    """
    Randomize A/B/C/D while preserving:
      - which option is correct
      - explanations
    """
    options = {
        "A": q.get("option_a", ""),
        "B": q.get("option_b", ""),
        "C": q.get("option_c", ""),
        "D": q.get("option_d", ""),
    }
    expl = {
        "A": q.get("option_a_explanation", ""),
        "B": q.get("option_b_explanation", ""),
        "C": q.get("option_c_explanation", ""),
        "D": q.get("option_d_explanation", ""),
    }
    correct = q.get("correct_answer", "A").strip().upper() or "A"

    old_labels = ["A", "B", "C", "D"]
    random.shuffle(old_labels)  # permutation of original labels

    new_options = {}
    new_expl = {}
    new_correct = None
    new_labels = ["A", "B", "C", "D"]

    for new_label, old_label in zip(new_labels, old_labels):
        new_options[new_label] = options[old_label]
        new_expl[new_label] = expl[old_label]
        if old_label == correct:
            new_correct = new_label

    # Update q in-place-like
    q["option_a"] = new_options["A"]
    q["option_b"] = new_options["B"]
    q["option_c"] = new_options["C"]
    q["option_d"] = new_options["D"]

    q["option_a_explanation"] = new_expl["A"]
    q["option_b_explanation"] = new_expl["B"]
    q["option_c_explanation"] = new_expl["C"]
    q["option_d_explanation"] = new_expl["D"]

    q["correct_answer"] = new_correct or "A"
    return q

# ==========================
# CLAUDE EVALUATION
# ==========================

def build_eval_prompt(
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    topic: str,
    subtopic: str
) -> str:
    return f"""
You are evaluating the quality of the following Java multiple-choice question.

Question: \"\"\"{question}\"\"\"
Options:
A) {option_a}
B) {option_b}
C) {option_c}
D) {option_d}

Correct Answer: {answer}

Option A explanation: {a_explanation}
Option B explanation: {b_explanation}
Option C explanation: {c_explanation}
Option D explanation: {d_explanation}
Overall Explanation: {main_explanation}

Bloom level: "{bloom_level}"
Topic: "{topic}"
Subtopic: "{subtopic}"

Rate each dimension from 1 (very poor) to 5 (excellent):

1. Relevance — how well it assesses the intended concept.
2. Bloom Alignment — how well it matches the given Bloom level.
3. Accuracy — technical correctness and clarity.
4. Explainability — clarity and helpfulness of explanations.

Return your scores in this exact format (no extra commentary):

relevance=X,bloom=Y,accuracy=Z,explainability=W,"Justification: [short text]"
"""

def parse_eval_response(response: str) -> Optional[Dict[str, Any]]:
    pattern = r"relevance[=:]\s*(\d).*?bloom[=:]\s*(\d).*?accuracy[=:]\s*(\d).*?explainability[=:]\s*(\d).*?\"Justification:\s*(.+?)\"$"
    match = re.search(pattern, response, re.DOTALL | re.IGNORECASE)
    if not match:
        print("⚠️ Could not parse eval response:\n", response)
        return None

    try:
        return {
            "eval_relevance": int(match.group(1)),
            "eval_bloom_alignment": int(match.group(2)),
            "eval_accuracy": int(match.group(3)),
            "eval_explainability": int(match.group(4)),
            "eval_justification": match.group(5).strip(),
        }
    except Exception as e:
        print("⚠️ Error parsing eval scores:", e)
        return None

def evaluate_with_claude(
    anthropic_client: Anthropic,
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    topic: str,
    subtopic: str
) -> Optional[Dict[str, Any]]:
    prompt = build_eval_prompt(
        question,
        option_a,
        option_b,
        option_c,
        option_d,
        answer,
        a_explanation,
        b_explanation,
        c_explanation,
        d_explanation,
        main_explanation,
        bloom_level,
        topic,
        subtopic
    )
    try:
        resp = anthropic_client.messages.create(
            model=CLAUDE_MODEL,
            max_tokens=300,
            temperature=0.2,
            messages=[{"role": "user", "content": prompt}],
        )
        text = resp.content[0].text.strip()
        return parse_eval_response(text)
    except Exception as e:
        print("⚠️ Claude evaluation failed:", e)
        return None

# ==========================
# CSV HELPERS
# ==========================

def ensure_output_csv(path: str):
    if not os.path.exists(path):
        with open(path, "w", newline="", encoding="utf-8") as f:
            writer = csv.writer(f)
            writer.writerow(CSV_COLUMNS)

def append_row_to_csv(path: str, row: Dict[str, Any]):
    with open(path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=CSV_COLUMNS)
        writer.writerow(row)

# ==========================
# MAIN PIPELINE FOR ONE ITEM
# ==========================

def run_pipeline_for_pair(
    openai_client: OpenAI,
    anthropic_client: Anthropic,
    slide_chunks: List[SlideChunk],
    slide_embeddings: List[np.ndarray],
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str,
    csv_path: str
) -> Optional[Dict[str, Any]]:
    # Build query for scoping
    query = f"{topic} - {subtopic}"
    slide_files, indices = get_relevant_slide_scope(
        openai_client,
        slide_chunks,
        slide_embeddings,
        query,
        k=4
    )

    if DEBUG:
        print(f"\nDEBUG: Topic '{topic}' | Subtopic '{subtopic}' | Bloom '{bloom_level}'")
        if slide_files:
            print("DEBUG: Using slide files:")
            for f in slide_files:
                print("  -", f)
        else:
            print("DEBUG: No slide files strongly matched this subtopic.")
        print()

    # Build creative, bounded prompt (no slide text)
    prompt = build_prompt(topic, subtopic, bloom_level, subject)

    # Generate question
    raw_response = generate_with_openai(openai_client, prompt)
    q = parse_question_response(raw_response)

    required = ["question", "option_a", "option_b", "option_c", "option_d", "correct_answer"]
    if not all(q.get(k, "").strip() for k in required):
        print("❌ Incomplete question structure, skipping.")
        return None

    # Randomize options (correct-answer position) in code
    q = shuffle_question_structure(q)

    # Evaluate with Claude (on shuffled item)
    eval_scores = evaluate_with_claude(
        anthropic_client,
        question=q["question"],
        option_a=q["option_a"],
        option_b=q["option_b"],
        option_c=q["option_c"],
        option_d=q["option_d"],
        answer=q["correct_answer"],
        a_explanation=q.get("option_a_explanation", ""),
        b_explanation=q.get("option_b_explanation", ""),
        c_explanation=q.get("option_c_explanation", ""),
        d_explanation=q.get("option_d_explanation", ""),
        main_explanation=q.get("main_explanation", ""),
        bloom_level=bloom_level,
        topic=topic,
        subtopic=subtopic,
    )

    # Pack slide scope metadata (no slide text)
    chunk_meta = [
        {"file": slide_chunks[i].slide_file, "slide_index": slide_chunks[i].slide_index}
        for i in indices
    ]
    retrieved_files = sorted(set(slide_files))
    retrieved_chunk_indices = json.dumps(chunk_meta, ensure_ascii=False)

    # Build row
    row = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.now().isoformat(),
        "topic": topic,
        "subtopic": subtopic,
        "subject": subject,
        "bloom_level": bloom_level,
        "cognitive_process": "",  # could derive from Bloom if desired
        "kc_tags": "; ".join(q.get("tags_list", [])),
        "question_stem": q["question"],
        "option_a": q["option_a"],
        "option_b": q["option_b"],
        "option_c": q["option_c"],
        "option_d": q["option_d"],
        "correct_answer": q["correct_answer"],
        "a_explanation": q.get("option_a_explanation", ""),
        "b_explanation": q.get("option_b_explanation", ""),
        "c_explanation": q.get("option_c_explanation", ""),
        "d_explanation": q.get("option_d_explanation", ""),
        "main_explanation": q.get("main_explanation", ""),
        "item_type": "MCQ",
        "eval_relevance": None,
        "eval_bloom_alignment": None,
        "eval_accuracy": None,
        "eval_explainability": None,
        "eval_justification": "",
        "raw_model_response": raw_response,
        "retrieved_slide_files": "; ".join(retrieved_files),
        "retrieved_chunk_indices": retrieved_chunk_indices,
        "generation_seed_prompt": prompt,
    }

    if eval_scores:
        row.update(eval_scores)

    append_row_to_csv(csv_path, row)
    return row

# ==========================
# MAIN DRIVER
# ==========================

def main():
    openai_client, anthropic_client = init_clients()

    # Load slides + embeddings once
    slide_chunks, _ = load_slide_chunks()
    texts = [c.text for c in slide_chunks]
    slide_embeddings = embed_texts(openai_client, texts) if texts else []

    ensure_output_csv(OUTPUT_CSV)

    questions_generated = 0

    for topic, subtopics in TOPIC_MAP.items():
        for subtopic in subtopics:
            for bloom in BLOOM_LEVELS:
                for i in range(QUESTIONS_PER_PAIR):
                    print(f"\n>> {topic} | {subtopic} | {bloom} | Q{i+1}")
                    try:
                        row = run_pipeline_for_pair(
                            openai_client,
                            anthropic_client,
                            slide_chunks,
                            slide_embeddings,
                            topic,
                            subtopic,
                            bloom,
                            SUBJECT,
                            OUTPUT_CSV
                        )
                        if row:
                            questions_generated += 1
                            print("✅ Question generated and saved.")
                        else:
                            print("❌ Generation failed for this question.")
                    except Exception as e:
                        print("🔥 Critical error:", e)

                    time.sleep(1.5)  # be nice to the APIs

    print("\n🎉 Done! Total questions generated:", questions_generated)
    print("📄 Saved to:", OUTPUT_CSV)


if __name__ == "__main__":
    main()






































# """
# question_generation.py

# Hybrid LLM pipeline for generating Java MCQs + adaptive metadata,
# restricted to the scope of your lecture slides.

# - Uses OpenAI (GPT-4.1 by default) to generate questions + metadata
# - Uses Anthropic Claude Sonnet to refine / validate / score items
# - Restricts generation using text chunks extracted from .pptx slides
# - Writes results to `data/java_questions_adaptive.csv`

# Environment variables required:
#     OPENAI_API_KEY       - for OpenAI
#     ANTHROPIC_API_KEY    - for Anthropic Claude

# Directory layout (expected):
#     question_generation/
#         question_generation.py   <-- this file
#         slides/                  <-- put all your .pptx here
#         data/
#             java_questions_adaptive.csv  <-- created/extended by this script

# Run:
#     python question_generation.py
# """

# import os
# import csv
# import json
# import uuid
# import time
# from dataclasses import dataclass
# from typing import List, Dict, Any, Optional, Tuple

# import pandas as pd
# from pptx import Presentation

# from openai import OpenAI
# from anthropic import Anthropic

# # ==========================
# # CONFIG
# # ==========================

# OPENAI_MODEL = "gpt-4.1"
# CLAUDE_MODEL = "claude-haiku-4-5-20251001"


# SLIDES_DIR = os.path.join(os.path.dirname(__file__), "slides")
# DATA_DIR = os.path.join(os.path.dirname(__file__), "data")
# os.makedirs(DATA_DIR, exist_ok=True)

# OUTPUT_CSV = os.path.join(DATA_DIR, "java_questions_adaptive.csv")

# DEBUG = True          # minimal transparent debug info
# TOP_K_CHUNKS = 5      # number of slide chunks to include as context
# MAX_CONTEXT_CHARS = 4000

# BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]

# TOPIC_MAP = {
#     "Java Fundamentals": [
#         "Variables and data types",
#         "Operators and expressions",
#         "Basic syntax and structure",
#         "Wrapper classes and autoboxing"
#     ],
#     "Control Flow": [
#         "if/else branching",
#         "switch statements",
#         "boolean expressions"
#     ],
#     "Loops": [
#         "for loop",
#         "while loop",
#         "do-while loop",
#         "loop control (break/continue)"
#     ],
#     "Arrays": [
#         "1D arrays",
#         "Array iteration",
#         "Common array errors"
#     ],
#     "Strings": [
#         "String immutability",
#         "Common String methods",
#         "String comparison and interning"
#     ],
#     "Methods": [
#         "Parameter passing",
#         "Method signatures",
#         "Return values"
#     ],
#     "Objects and Classes": [
#         "Fields and methods",
#         "Constructors",
#         "Instance vs static members",
#         "toString and equals overrides"
#     ],
#     "Encapsulation": [
#         "Access modifiers",
#         "Getters and setters",
#         "Information hiding"
#     ],
#     "Inheritance": [
#         "Superclass/subclass relationships",
#         "Method overriding",
#         "super keyword"
#     ],
#     "Polymorphism": [
#         "Dynamic dispatch",
#         "Upcasting and downcasting",
#         "Method binding"
#     ],
#     "Abstract Classes": [
#         "Abstract methods",
#         "Partial implementation",
#         "Concrete subclasses of abstract classes"
#     ],
#     "Interfaces": [
#         "Interface contracts",
#         "Implementing multiple interfaces",
#         "Functional interfaces"
#     ],
#     "Generics": [
#         "Generic classes",
#         "Generic methods",
#         "Type parameters and type safety"
#     ],
#     "Collections": [
#         "Lists (e.g., ArrayList)",
#         "Sets and uniqueness",
#         "Maps (key-value pairs)",
#         "Iteration over collections"
#     ],
#     "JavaFX": [
#         "Scene graph",
#         "UI controls and layout",
#         "Stages and scenes"
#     ],
#     "Event-Driven Programming": [
#         "Event handlers",
#         "Listener patterns",
#         "JavaFX event handling with lambdas"
#     ]
# }

# # How many questions per (topic, subtopic, bloom) pair
# QUESTIONS_PER_PAIR = 3

# SUBJECT = "Java"

# # ==========================
# # LLM CLIENTS
# # ==========================

# def get_openai_client() -> OpenAI:
#     api_key = os.environ.get("OPENAI_API_KEY")
#     if not api_key:
#         raise RuntimeError("OPENAI_API_KEY environment variable is not set.")
#     return OpenAI(api_key=api_key)

# def get_claude_client() -> Anthropic:
#     api_key = os.environ.get("ANTHROPIC_API_KEY")
#     if not api_key:
#         raise RuntimeError("ANTHROPIC_API_KEY environment variable is not set.")
#     return Anthropic(api_key=api_key)

# # ==========================
# # SLIDE PROCESSING
# # ==========================

# @dataclass
# class SlideChunk:
#     slide_file: str
#     slide_index: int
#     text: str

# def extract_text_from_pptx(path: str) -> List[str]:
#     prs = Presentation(path)
#     texts = []
#     for i, slide in enumerate(prs.slides):
#         parts = []
#         for shape in slide.shapes:
#             if hasattr(shape, "text") and shape.has_text_frame:
#                 t = shape.text.strip()
#                 if t:
#                     parts.append(t)
#         slide_text = "\n".join(parts).strip()
#         if slide_text:
#             texts.append(slide_text)
#     return texts

# def load_slide_chunks(slides_dir: str) -> List[SlideChunk]:
#     chunks: List[SlideChunk] = []
#     if not os.path.isdir(slides_dir):
#         return chunks
#     for fname in os.listdir(slides_dir):
#         if not fname.lower().endswith(".pptx"):
#             continue
#         full = os.path.join(slides_dir, fname)
#         try:
#             slide_texts = extract_text_from_pptx(full)
#             for idx, txt in enumerate(slide_texts):
#                 chunks.append(SlideChunk(slide_file=fname, slide_index=idx, text=txt))
#         except Exception as e:
#             print(f"⚠️ Failed to read {full}: {e}")
#     if DEBUG:
#         print(f"Loaded {len(chunks)} slide chunks from {slides_dir}")
#     return chunks

# def simple_score(query: str, text: str) -> int:
#     query_tokens = {t.lower() for t in query.replace(",", " ").split() if len(t) > 2}
#     text_tokens = {t.lower() for t in text.replace(",", " ").split() if len(t) > 2}
#     return len(query_tokens & text_tokens)

# def build_context_for(topic: str,
#                       subtopic: str,
#                       slide_chunks: List[SlideChunk],
#                       top_k: int = TOP_K_CHUNKS,
#                       max_chars: int = MAX_CONTEXT_CHARS
#                       ) -> Tuple[str, List[Tuple[str, int]]]:
#     if not slide_chunks:
#         return "", []

#     query = f"{topic} - {subtopic}"
#     scored = []
#     for ch in slide_chunks:
#         s = simple_score(query, ch.text)
#         if s > 0:
#             scored.append((s, ch))
#     if not scored:
#         # if nothing matched, just take first N chunks
#         scored = [(1, ch) for ch in slide_chunks[:top_k]]

#     scored.sort(key=lambda x: x[0], reverse=True)
#     selected = [ch for _, ch in scored[:top_k]]

#     context_parts = []
#     debug_chunks = []
#     total_chars = 0
#     for ch in selected:
#         snippet = ch.text.strip()
#         if total_chars + len(snippet) > max_chars:
#             break
#         context_parts.append(snippet)
#         total_chars += len(snippet)
#         debug_chunks.append((ch.slide_file, ch.slide_index))

#     context = "\n\n---\n\n".join(context_parts)
#     return context, debug_chunks

# # ==========================
# # PROMPT BUILDING
# # ==========================

# def build_generation_prompt(topic: str,
#                             subtopic: str,
#                             bloom_level: str,
#                             subject: str,
#                             context: str) -> str:
#     return f"""
# You are an expert Java instructor and assessment designer.

# Your task is to generate ONE high-quality multiple-choice question (MCQ) for an undergraduate {subject} course.

# The question MUST:

# - Be about the topic: "{topic}"
# - Focus specifically on the subtopic: "{subtopic}"
# - Operate at Bloom level: "{bloom_level}"
# - Be fully contained within the scope of the provided course material (do NOT introduce concepts that clearly require content beyond the slides).
# - Be creative, varied, and realistic. Do NOT mention slides, PPT, lecture, or any phrasing like "according to the slides".

# Use the following COURSE CONTEXT (extracted from lecture slides) ONLY as background knowledge. You do NOT need to quote it. Just use it as a content boundary:

# [COURSE CONTEXT START]
# {context}
# [COURSE CONTEXT END]

# Return your answer as a SINGLE JSON object with the following structure (no extra commentary, no markdown):

# {{
#   "topic": "{topic}",
#   "subtopic": "{subtopic}",
#   "subject": "{subject}",
#   "bloom_level": "{bloom_level}",
#   "question_stem": "...",
#   "options": {{
#     "A": "...",
#     "B": "...",
#     "C": "...",
#     "D": "..."
#   }},
#   "correct_answer": "A",
#   "a_explanation": "...",
#   "b_explanation": "...",
#   "c_explanation": "...",
#   "d_explanation": "...",
#   "main_explanation": "...",

#   "cognitive_process": "e.g., recall / explain / apply / analyze / evaluate",
#   "kc_tags": ["short", "keyword", "tags"],
#   "item_type": "conceptual MCQ or code-based MCQ",

#   "predicted_difficulty_level": 3,
#   "predicted_difficulty_label": "easy / medium / hard",
#   "reasoning_depth": 2,
#   "linguistic_complexity": "simple / moderate / complex",
#   "estimated_time_seconds": 60,

#   "distractor_analysis": "Brief description of how distractors relate to common misunderstandings.",
#   "variant_group_id": "short_id_grouping_similar_items",

#   "irt_difficulty_b": 0.0,
#   "irt_discrimination_a": 1.0,
#   "irt_guessing_c": 0.25,

#   "remediation_reference": "Natural language hint on what to review.",
#   "safety_notes": "",
#   "bias_notes": "",

#   "misconception_tags_per_option": {{
#     "A": ["tag1_if_any"],
#     "B": ["tag2_if_any"],
#     "C": [],
#     "D": []
#   }}
# }}
# """

# def build_claude_refinement_prompt(raw_json: str,
#                                    topic: str,
#                                    subtopic: str,
#                                    bloom_level: str,
#                                    subject: str,
#                                    context: str) -> str:
#     return f"""
# You are reviewing a JSON-defined multiple-choice question for an undergraduate {subject} course.

# The question was generated to:
# - Be about topic: "{topic}"
# - Focus on subtopic: "{subtopic}"
# - Match Bloom level: "{bloom_level}"
# - Stay within the scope of the provided course material

# Here is the COURSE CONTEXT extracted from slides (this is the content boundary; do not introduce material clearly beyond this):

# [COURSE CONTEXT START]
# {context}
# [COURSE CONTEXT END]

# Here is the ORIGINAL JSON for the question:

# [QUESTION JSON START]
# {raw_json}
# [QUESTION JSON END]

# Your tasks:

# 1. Parse the JSON and check:
#    - Technical correctness of the Java content.
#    - That the correct option is truly correct and the distractors are plausible but wrong.
#    - That the question matches the intended Bloom level.
#    - That the content is within the likely scope of the course material.

# 2. If needed, minimally adjust:
#    - question_stem
#    - options and correct_answer
#    - explanations
#    - kc_tags
#    - predicted_difficulty_level / label
#    - misconception_tags_per_option
#    - any other metadata fields

# 3. Add an evaluation object with:
#    - eval_relevance: 1-5
#    - eval_bloom_alignment: 1-5
#    - eval_accuracy: 1-5
#    - eval_explainability: 1-5
#    - eval_justification: short justification string

# Return a SINGLE JSON object with the same top-level fields as the original question, plus an extra field:

# "evaluation": {{
#   "eval_relevance": 5,
#   "eval_bloom_alignment": 5,
#   "eval_accuracy": 5,
#   "eval_explainability": 5,
#   "eval_justification": "..."
# }}

# Output ONLY valid JSON. No markdown, no comments, no backticks.
# """

# # ==========================
# # LLM CALLS
# # ==========================

# def call_openai_generate(prompt: str, client: OpenAI) -> Dict[str, Any]:
#     resp = client.chat.completions.create(
#         model=OPENAI_MODEL,
#         messages=[
#             {"role": "system", "content": "You are a careful Java tutor and assessment designer. Always respond with strict JSON only."},
#             {"role": "user", "content": prompt}
#         ],
#         temperature=0.8,
#         max_tokens=800,
#     )
#     text = resp.choices[0].message.content.strip()
#     try:
#         return json.loads(text)
#     except Exception as e:
#         raise ValueError(f"Failed to parse OpenAI JSON: {e}\nRaw text:\n{text}")

# def call_claude_refine(raw_item: Dict[str, Any],
#                        topic: str,
#                        subtopic: str,
#                        bloom_level: str,
#                        subject: str,
#                        context: str,
#                        client: Anthropic) -> Dict[str, Any]:
#     raw_json = json.dumps(raw_item, ensure_ascii=False, indent=2)
#     prompt = build_claude_refinement_prompt(
#         raw_json=raw_json,
#         topic=topic,
#         subtopic=subtopic,
#         bloom_level=bloom_level,
#         subject=subject,
#         context=context
#     )
#     resp = client.messages.create(
#         model=CLAUDE_MODEL,
#         max_tokens=900,
#         temperature=0.2,
#         messages=[{"role": "user", "content": prompt}]
#     )
#     text = resp.content[0].text.strip()
#     try:
#         return json.loads(text)
#     except Exception as e:
#         raise ValueError(f"Failed to parse Claude JSON: {e}\nRaw text:\n{text}")

# # ==========================
# # CSV HANDLING
# # ==========================

# CSV_COLUMNS = [
#     "id",
#     "timestamp",
#     "topic",
#     "subtopic",
#     "subject",
#     "bloom_level",
#     "cognitive_process",
#     "kc_tags",
#     "question_stem",
#     "option_a",
#     "option_b",
#     "option_c",
#     "option_d",
#     "correct_answer",
#     "a_explanation",
#     "b_explanation",
#     "c_explanation",
#     "d_explanation",
#     "main_explanation",
#     "item_type",
#     "predicted_difficulty_level",
#     "predicted_difficulty_label",
#     "reasoning_depth",
#     "linguistic_complexity",
#     "estimated_time_seconds",
#     "distractor_analysis",
#     "variant_group_id",
#     "irt_difficulty_b",
#     "irt_discrimination_a",
#     "irt_guessing_c",
#     "remediation_reference",
#     "safety_notes",
#     "bias_notes",
#     "raw_model_response",
#     "retrieved_slide_files",
#     "retrieved_chunk_indices",
#     "generation_seed_prompt",
#     "eval_relevance",
#     "eval_bloom_alignment",
#     "eval_accuracy",
#     "eval_explainability",
#     "eval_justification",
#     "misconception_tags_per_option"
# ]

# def ensure_output_csv(path: str):
#     if not os.path.exists(path):
#         os.makedirs(os.path.dirname(path), exist_ok=True)
#         with open(path, "w", newline="", encoding="utf-8") as f:
#             writer = csv.writer(f)
#             writer.writerow(CSV_COLUMNS)

# def append_row_to_csv(path: str, row: Dict[str, Any]):
#     with open(path, "a", newline="", encoding="utf-8") as f:
#         writer = csv.DictWriter(f, fieldnames=CSV_COLUMNS)
#         writer.writerow(row)

# # ==========================
# # PIPELINE
# # ==========================

# def generate_one_item(topic: str,
#                       subtopic: str,
#                       bloom_level: str,
#                       subject: str,
#                       slide_chunks: List[SlideChunk],
#                       openai_client: OpenAI,
#                       claude_client: Anthropic) -> Optional[pd.Series]:
#     context, dbg_chunks = build_context_for(topic, subtopic, slide_chunks)

#     if DEBUG:
#         print(f"\n--- Generating for Topic: {topic} | Subtopic: {subtopic} | Bloom: {bloom_level}")
#         print("Slide chunks used:")
#         for fname, idx in dbg_chunks:
#             print(f"  - {fname} (slide {idx})")

#     gen_prompt = build_generation_prompt(
#         topic=topic,
#         subtopic=subtopic,
#         bloom_level=bloom_level,
#         subject=subject,
#         context=context
#     )

#     try:
#         raw_item = call_openai_generate(gen_prompt, openai_client)
#     except Exception as e:
#         print(f"❌ OpenAI generation failed: {e}")
#         return None

#     try:
#         refined_item = call_claude_refine(
#             raw_item=raw_item,
#             topic=topic,
#             subtopic=subtopic,
#             bloom_level=bloom_level,
#             subject=subject,
#             context=context,
#             client=claude_client
#         )
#     except Exception as e:
#         print(f"❌ Claude refinement failed, falling back to raw item: {e}")
#         refined_item = raw_item
#         refined_item.setdefault("evaluation", {
#             "eval_relevance": None,
#             "eval_bloom_alignment": None,
#             "eval_accuracy": None,
#             "eval_explainability": None,
#             "eval_justification": "Claude refinement failed; using raw item."
#         })

#     evaluation = refined_item.get("evaluation", {})

#     options = refined_item.get("options", {})
#     mc_tags = refined_item.get("misconception_tags_per_option", {})

#     row = {
#         "id": str(uuid.uuid4()),
#         "timestamp": pd.Timestamp.now().isoformat(),
#         "topic": topic,
#         "subtopic": subtopic,
#         "subject": subject,
#         "bloom_level": refined_item.get("bloom_level", bloom_level),
#         "cognitive_process": refined_item.get("cognitive_process", ""),
#         "kc_tags": "; ".join(refined_item.get("kc_tags", [])),
#         "question_stem": refined_item.get("question_stem", ""),
#         "option_a": options.get("A", ""),
#         "option_b": options.get("B", ""),
#         "option_c": options.get("C", ""),
#         "option_d": options.get("D", ""),
#         "correct_answer": refined_item.get("correct_answer", ""),
#         "a_explanation": refined_item.get("a_explanation", ""),
#         "b_explanation": refined_item.get("b_explanation", ""),
#         "c_explanation": refined_item.get("c_explanation", ""),
#         "d_explanation": refined_item.get("d_explanation", ""),
#         "main_explanation": refined_item.get("main_explanation", ""),
#         "item_type": refined_item.get("item_type", ""),
#         "predicted_difficulty_level": refined_item.get("predicted_difficulty_level", None),
#         "predicted_difficulty_label": refined_item.get("predicted_difficulty_label", ""),
#         "reasoning_depth": refined_item.get("reasoning_depth", None),
#         "linguistic_complexity": refined_item.get("linguistic_complexity", ""),
#         "estimated_time_seconds": refined_item.get("estimated_time_seconds", None),
#         "distractor_analysis": refined_item.get("distractor_analysis", ""),
#         "variant_group_id": refined_item.get("variant_group_id", ""),
#         "irt_difficulty_b": refined_item.get("irt_difficulty_b", None),
#         "irt_discrimination_a": refined_item.get("irt_discrimination_a", None),
#         "irt_guessing_c": refined_item.get("irt_guessing_c", None),
#         "remediation_reference": refined_item.get("remediation_reference", ""),
#         "safety_notes": refined_item.get("safety_notes", ""),
#         "bias_notes": refined_item.get("bias_notes", ""),
#         "raw_model_response": json.dumps(raw_item, ensure_ascii=False),
#         "retrieved_slide_files": "; ".join(sorted({f for f, _ in dbg_chunks})),
#         "retrieved_chunk_indices": json.dumps(
#             [{"file": f, "slide_index": i} for f, i in dbg_chunks],
#             ensure_ascii=False
#         ),
#         "generation_seed_prompt": gen_prompt,
#         "eval_relevance": evaluation.get("eval_relevance"),
#         "eval_bloom_alignment": evaluation.get("eval_bloom_alignment"),
#         "eval_accuracy": evaluation.get("eval_accuracy"),
#         "eval_explainability": evaluation.get("eval_explainability"),
#         "eval_justification": evaluation.get("eval_justification", ""),
#         "misconception_tags_per_option": json.dumps(mc_tags, ensure_ascii=False)
#     }

#     return pd.Series(row)

# def run_generation():
#     print("🔧 Loading slide chunks...")
#     slide_chunks = load_slide_chunks(SLIDES_DIR)
#     if not slide_chunks:
#         print(f"⚠️ No slides found in {SLIDES_DIR}. The model will still run, but without slide restriction.")

#     ensure_output_csv(OUTPUT_CSV)

#     openai_client = get_openai_client()
#     claude_client = get_claude_client()

#     total_generated = 0

#     for topic, subtopics in TOPIC_MAP.items():
#         for subtopic in subtopics:
#             for bloom_level in BLOOM_LEVELS:
#                 for i in range(QUESTIONS_PER_PAIR):
#                     print(f"\n=== {topic} | {subtopic} | {bloom_level} | Q#{i+1} ===")
#                     try:
#                         series = generate_one_item(
#                             topic=topic,
#                             subtopic=subtopic,
#                             bloom_level=bloom_level,
#                             subject=SUBJECT,
#                             slide_chunks=slide_chunks,
#                             openai_client=openai_client,
#                             claude_client=claude_client
#                         )
#                         if series is not None:
#                             append_row_to_csv(OUTPUT_CSV, series.to_dict())
#                             total_generated += 1
#                             print(f"✅ Saved question {total_generated} to CSV.")
#                         else:
#                             print("❌ Skipped due to generation failure.")
#                     except Exception as e:
#                         print(f"🔥 Critical error on {topic} / {subtopic} / {bloom_level}: {e}")
#                     time.sleep(1.0)  # gentle pacing

#     print(f"\n🎉 Generation complete! Total questions generated this run: {total_generated}")
#     print(f"Output CSV: {OUTPUT_CSV}")

# if __name__ == "__main__":
#     run_generation()
