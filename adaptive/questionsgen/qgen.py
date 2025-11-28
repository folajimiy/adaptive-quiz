#!/usr/bin/env python3
"""
slide_question_gen.py

Slide-bounded, creative Java MCQ generator with topics + subtopics.

- Automatically loads all .pptx slides from a ./slides folder
- Uses slide text only to build embeddings and map topics/subtopics to slide files
- Does NOT show slide text to the LLM (Option A: slide-bounded, not slide-anchored)
- Generates creative MCQs using OpenAI (GPT-4.1 family)
- Evaluates them using Claude (Sonnet)
- Logs topic, subtopic, Bloom level, and slide filenames used into a CSV

Requirements:
    pip install python-pptx openai anthropic numpy pandas

Environment variables:
    OPENAI_API_KEY
    CLAUDE_ANTHROPIC_KEY
"""

import os
import re
import csv
import uuid
import time
from datetime import datetime

import numpy as np
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from anthropic import Anthropic

# ==========================
# GLOBAL CONFIG
# ==========================

# Minimal debug: just show which slide files are used per (topic, subtopic)
DEBUG = True

# Models (you can change these strings if needed)
# OpenAI + Anthropic models
OPENAI_CHAT_MODEL = "gpt-4.1"
OPENAI_EMBED_MODEL = "text-embedding-3-small"
CLAUDE_MODEL = "claude-haiku-4-5-20251001"


# Output CSV and generation parameters
OUTPUT_CSV = "java_questions_slide_scoped_topics_subtopics.csv"
QUESTIONS_PER_CELL = 3  # questions per (topic, subtopic, Bloom level)

# Root paths
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SLIDES_FOLDER = os.path.join(SCRIPT_DIR, "slides")

# Subject label
SUBJECT = "Java"

# ==========================
# TOPICS + SUBTOPICS
# ==========================

TOPIC_MAP = {
    "Java Fundamentals": [
        "Variables and data types",
        "Operators and expressions",
        "Basic syntax and structure",
        "Wrapper classes and autoboxing"
    ],
    "Control Flow": [
        "if/else branching",
        "switch statements",
        "boolean expressions"
    ],
    "Loops": [
        "for loop",
        "while loop",
        "do-while loop",
        "loop control (break/continue)"
    ],
    "Arrays": [
        "1D arrays",
        "Array iteration",
        "Common array errors"
    ],
    "Strings": [
        "String immutability",
        "Common String methods",
        "String comparison and interning"
    ],
    "Methods": [
        "Parameter passing",
        "Method signatures",
        "Return values"
    ],
    "Objects and Classes": [
        "Fields and methods",
        "Constructors",
        "Instance vs static members",
        "toString and equals overrides"
    ],
    "Encapsulation": [
        "Access modifiers",
        "Getters and setters",
        "Information hiding"
    ],
    "Inheritance": [
        "Superclass/subclass relationships",
        "Method overriding",
        "super keyword"
    ],
    "Polymorphism": [
        "Dynamic dispatch",
        "Upcasting and downcasting",
        "Method binding"
    ],
    "Abstract Classes": [
        "Abstract methods",
        "Partial implementation",
        "Concrete subclasses of abstract classes"
    ],
    "Interfaces": [
        "Interface contracts",
        "Implementing multiple interfaces",
        "Functional interfaces"
    ],
    "Generics": [
        "Generic classes",
        "Generic methods",
        "Type parameters and type safety"
    ],
    "Collections": [
        "Lists (e.g., ArrayList)",
        "Sets and uniqueness",
        "Maps (key-value pairs)",
        "Iteration over collections"
    ],
    "JavaFX": [
        "Scene graph",
        "UI controls and layout",
        "Stages and scenes"
    ],
    "Event-Driven Programming": [
        "Event handlers",
        "Listener patterns",
        "JavaFX event handling with lambdas"
    ]
}

# ==========================
# API CLIENTS
# ==========================

def init_clients():
    openai_key = os.getenv("OPENAI_API_KEY")
    claude_key = os.getenv("CLAUDE_ANTHROPIC_KEY")

    if not openai_key:
        raise RuntimeError("OPENAI_API_KEY environment variable is not set.")
    if not claude_key:
        raise RuntimeError("CLAUDE_ANTHROPIC_KEY environment variable is not set.")

    openai_client = OpenAI(api_key=openai_key)
    anthropic_client = Anthropic(api_key=claude_key)
    return openai_client, anthropic_client

# ==========================
# SLIDE DISCOVERY & TEXT EXTRACTION
# ==========================

def get_all_pptx_files(folder: str):
    if not os.path.exists(folder):
        raise RuntimeError(f"Slides folder does not exist: {folder}")
    files = [
        os.path.join(folder, f)
        for f in os.listdir(folder)
        if f.lower().endswith(".pptx")
    ]
    if not files:
        raise RuntimeError(f"No .pptx slide files found in {folder}")
    return sorted(files)

def ppt_to_text(path: str) -> str:
    """Extract visible text from a PPTX file."""
    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)

# ==========================
# CHUNKING & EMBEDDINGS
# ==========================

def chunk_text(text: str, max_len: int = 1200, overlap: int = 200):
    """
    Simple character-wise chunker with overlap,
    preferring to break at newline boundaries.
    """
    text = text.strip()
    if not text:
        return []

    chunks = []
    n = len(text)
    start = 0

    while start < n:
        end = min(n, start + max_len)
        chunk = text[start:end]

        # Try to break at a reasonable newline
        last_break = max(chunk.rfind("\n\n"), chunk.rfind("\n"))
        if last_break > 200:  # don't break too early
            end = start + last_break
            chunk = text[start:end]

        chunk = chunk.strip()
        if len(chunk) > 50:
            chunks.append(chunk)

        start = max(start + 1, end - overlap)

    return chunks

def chunk_slides_with_sources(slide_paths):
    """
    For each slide deck, chunk its text and record which file
    each chunk came from.

    Returns:
        slide_chunks: list[str]
        chunk_sources: list[str]  (basename of PPTX for each chunk)
    """
    all_chunks = []
    all_sources = []

    for slide_path in slide_paths:
        base = os.path.basename(slide_path)
        print(f"📥 Loading & chunking slides from: {base}")
        text = ppt_to_text(slide_path)
        if not text.strip():
            continue
        slide_chunks = chunk_text(text, max_len=1200, overlap=200)
        all_chunks.extend(slide_chunks)
        all_sources.extend([base] * len(slide_chunks))

    print(f"📌 Total chunks: {len(all_chunks)}")
    return all_chunks, all_sources

def embed_texts(openai_client: OpenAI, texts):
    """Get embeddings for a list of texts using OpenAI."""
    if not texts:
        return []
    resp = openai_client.embeddings.create(
        model=OPENAI_EMBED_MODEL,
        input=texts
    )
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]

def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)

# ==========================
# RETRIEVAL (SCOPE ONLY)
# ==========================

def get_relevant_slide_scope(
    openai_client: OpenAI,
    query: str,
    slide_chunks,
    slide_embeddings,
    chunk_sources,
    k: int = 4
):
    """
    Uses embeddings to see which slide chunks are most related to the
    (topic + subtopic) query. Returns only slide filenames and indices
    — NO slide text is ever shown to the LLM.

    Returns:
        unique_files: sorted list of slide filenames involved
        indices: list of chunk indices chosen
    """
    if not slide_chunks or not slide_embeddings:
        return [], []

    query_emb = embed_texts(openai_client, [query])[0]
    sims = [cosine_similarity(query_emb, emb) for emb in slide_embeddings]
    top_indices = np.argsort(sims)[-k:][::-1]

    selected_files = [chunk_sources[i] for i in top_indices]
    unique_files = sorted(set(selected_files))

    return unique_files, list(top_indices)

# ==========================
# PROMPT BUILDING (CREATIVE, BOUNDED, NO SLIDE TEXT)
# ==========================

def build_prompt(
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str
) -> str:
    """
    Creative MCQ generation bounded by topic + subtopic and Bloom level.
    Slides are NOT mentioned and slide text is NOT included.
    """
    # A bit of extra safety: mention that it's for an intro/intermediate OOP class.
    return f"""
You are generating a high-quality, original multiple-choice question
for a university-level {subject} course focused on object-oriented programming.

You MUST follow these rules:

1. Stay strictly within the conceptual boundaries of:
   - Topic: "{topic}"
   - Subtopic: "{subtopic}"

2. Do NOT mention slides, lectures, or course materials.
3. Do NOT introduce advanced Java topics beyond standard introductory/
   intermediate OOP material (avoid networking, reflection, annotations,
   concurrency, advanced frameworks, etc.).
4. The question must be creative, non-trivial, and the stem MUST be ≤ 50 words.
5. Target Bloom level: "{bloom_level}".
6. There must be exactly ONE correct option.

Generate the MCQ using EXACTLY this format:

Question: [your question stem here]

Options:
A) [Option A text]
B) [Option B text]
C) [Option C text]
D) [Option D text]

Correct Answer: [A, B, C, or D]

Option A explanation: [why A is correct or incorrect]
Option B explanation: [why B is correct or incorrect]
Option C explanation: [why C is correct or incorrect]
Option D explanation: [why D is correct or incorrect]

Main Explanation: [brief overall explanation of the answer]

Tags: {topic}, {subtopic}, {bloom_level}
"""

# ==========================
# GENERATION & PARSING
# ==========================

def generate_with_openai(openai_client: OpenAI, prompt: str) -> str:
    resp = openai_client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a Java exam question generator. "
                    "You must strictly follow the response format and rules."
                )
            },
            {"role": "user", "content": prompt},
        ],
        temperature=0.9,
        top_p=0.95,
        max_tokens=500
    )
    return resp.choices[0].message.content.strip()

def parse_question_response(response: str) -> dict:
    patterns = {
        'question': r"[Qq]uestion:\s*(.+?)\s*[Oo]ptions?:",
        'option_a': r"A\)\s*(.+?)(?=\n[B-D]\)|\nCorrect Answer:|$)",
        'option_b': r"B\)\s*(.+?)(?=\n[C-D]\)|\nCorrect Answer:|$)",
        'option_c': r"C\)\s*(.+?)(?=\nD\)|\nCorrect Answer:|$)",
        'option_d': r"D\)\s*(.+?)(?=\nCorrect Answer:|$)",
        'correct_answer': r"Correct Answer:\s*([A-D])",
        'option_a_explanation': r"Option A explanation:\s*(.+?)(?=\nOption B explanation:|\nMain Explanation:|$)",
        'option_b_explanation': r"Option B explanation:\s*(.+?)(?=\nOption C explanation:|\nMain Explanation:|$)",
        'option_c_explanation': r"Option C explanation:\s*(.+?)(?=\nOption D explanation:|\nMain Explanation:|$)",
        'option_d_explanation': r"Option D explanation:\s*(.+?)(?=\nMain Explanation:|$)",
        'main_explanation': r"Main Explanation:\s*(.+?)(?=\nTags:|$)",
        'tags': r"Tags:\s*(.+)",
    }

    extracted = {}
    for key, pattern in patterns.items():
        m = re.search(pattern, response, re.DOTALL)
        extracted[key] = m.group(1).strip() if m else ""
    return extracted

# ==========================
# CLAUDE EVALUATION
# ==========================

def build_eval_prompt(
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    topic: str,
    subtopic: str
) -> str:
    return f"""
You are evaluating a Java multiple-choice question for a course on object-oriented programming.

Question: \"\"\"{question}\"\"\"
Options:
A) {option_a}
B) {option_b}
C) {option_c}
D) {option_d}

Correct Answer: {answer}

Option A explanation: {a_explanation}
Option B explanation: {b_explanation}
Option C explanation: {c_explanation}
Option D explanation: {d_explanation}
Overall Explanation: {main_explanation}

The question is intended to assess:
- Topic: "{topic}"
- Subtopic: "{subtopic}"
- Bloom level: "{bloom_level}"

Rate the question on a scale of 1 (very poor) to 5 (excellent) on:
1. Relevance — Does it directly assess the intended topic/subtopic?
2. Accuracy — Are the correct answer and distractors technically accurate and unambiguous?
3. Explainability — Are explanations clear, specific, and educational?
4. Bloom Alignment — Does the cognitive demand match Bloom: "{bloom_level}"?

Return ONLY in this exact format (no extra commentary):

relevance=X,bloom=Y,accuracy=Z,explainability=W,"Justification: [text]"
"""

def parse_eval_response(response: str):
    pattern = (
        r"relevance[=:]\s*(\d).*?"
        r"bloom[=:]\s*(\d).*?"
        r"accuracy[=:]\s*(\d).*?"
        r"explainability[=:]\s*(\d).*?"
        r"Justification:\s*\"?(.+?)\"?$"
    )
    m = re.search(pattern, response, re.DOTALL | re.IGNORECASE)
    if not m:
        print("⚠️ Could not parse evaluation response:")
        print(response)
        return None

    try:
        return {
            "eval_relevance": int(m.group(1)),
            "eval_bloom_alignment": int(m.group(2)),
            "eval_accuracy": int(m.group(3)),
            "eval_explainability": int(m.group(4)),
            "eval_justification": m.group(5).strip(),
        }
    except Exception as e:
        print("⚠️ Error parsing eval scores:", e)
        return None

def evaluate_with_claude(
    anthropic_client: Anthropic,
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    topic: str,
    subtopic: str
):
    prompt = build_eval_prompt(
        question, option_a, option_b, option_c, option_d, answer,
        a_explanation, b_explanation, c_explanation, d_explanation,
        main_explanation, bloom_level, topic, subtopic
    )
    resp = anthropic_client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=300,
        temperature=0.2,
        messages=[{"role": "user", "content": prompt}],
    )
    text = resp.content[0].text
    return parse_eval_response(text)

# ==========================
# PIPELINE
# ==========================

def run_pipeline(
    openai_client: OpenAI,
    anthropic_client: Anthropic,
    slide_chunks,
    slide_embeddings,
    chunk_sources,
    topic: str,
    subtopic: str,
    bloom_level: str,
    subject: str,
    csv_path: str
):
    # Build query for scoping (topic + subtopic)
    query = f"{topic} - {subtopic}"
    slide_files, indices = get_relevant_slide_scope(
        openai_client,
        query,
        slide_chunks,
        slide_embeddings,
        chunk_sources,
        k=4
    )

    if DEBUG:
        print(f"\nDEBUG: Topic '{topic}' | Subtopic '{subtopic}'")
        if slide_files:
            print("DEBUG: Using slide files:")
            for f in slide_files:
                print("  -", f)
        else:
            print("DEBUG: No slide files strongly matched this subtopic.")
        print()

    # Build creative, bounded prompt (no slide text)
    prompt = build_prompt(topic, subtopic, bloom_level, subject)

    # Generate question
    raw_response = generate_with_openai(openai_client, prompt)
    q = parse_question_response(raw_response)

    required = ["question", "option_a", "option_b", "option_c", "option_d", "correct_answer"]
    if not all(q.get(k, "").strip() for k in required):
        print("❌ Incomplete question structure, skipping.")
        return None

    # Evaluate with Claude
    eval_scores = evaluate_with_claude(
        anthropic_client,
        question=q["question"],
        option_a=q["option_a"],
        option_b=q["option_b"],
        option_c=q["option_c"],
        option_d=q["option_d"],
        answer=q["correct_answer"],
        a_explanation=q.get("option_a_explanation", ""),
        b_explanation=q.get("option_b_explanation", ""),
        c_explanation=q.get("option_c_explanation", ""),
        d_explanation=q.get("option_d_explanation", ""),
        main_explanation=q.get("main_explanation", ""),
        bloom_level=bloom_level,
        topic=topic,
        subtopic=subtopic,
    )

    row = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.now().isoformat(),
        "topic": topic,
        "subtopic": subtopic,
        "bloom_level": bloom_level,
        "question_stem": q["question"],
        "option_a": q["option_a"],
        "option_b": q["option_b"],
        "option_c": q["option_c"],
        "option_d": q["option_d"],
        "correct_answer": q["correct_answer"],
        "a_explanation": q.get("option_a_explanation", ""),
        "b_explanation": q.get("option_b_explanation", ""),
        "c_explanation": q.get("option_c_explanation", ""),
        "d_explanation": q.get("option_d_explanation", ""),
        "main_explanation": q.get("main_explanation", ""),
        "raw_model_response": raw_response,
        "retrieved_slide_files": ", ".join(slide_files),
        "retrieved_chunk_indices": ", ".join(str(i) for i in indices),
    }

    if eval_scores:
        row.update(eval_scores)
    else:
        row.update({
            "eval_relevance": None,
            "eval_bloom_alignment": None,
            "eval_accuracy": None,
            "eval_explainability": None,
            "eval_justification": None,
        })

    # Append to CSV
    file_exists = os.path.exists(csv_path)
    with open(csv_path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=row.keys())
        if not file_exists:
            writer.writeheader()
        writer.writerow(row)

    return row

# ==========================
# MAIN
# ==========================

def main():
    openai_client, anthropic_client = init_clients()

    # Discover slide files
    slide_paths = get_all_pptx_files(SLIDES_FOLDER)
    print("📄 Found slide decks:")
    for p in slide_paths:
        print("   ", os.path.basename(p))

    # Chunk slides and embed
    slide_chunks, chunk_sources = chunk_slides_with_sources(slide_paths)
    print("🔍 Embedding slide chunks (one-time)...")
    slide_embeddings = embed_texts(openai_client, slide_chunks)
    print("✅ Embeddings ready.")

    print(f"\n💾 Output CSV: {OUTPUT_CSV}")
    total_questions = 0

    for topic, subtopics in TOPIC_MAP.items():
        for subtopic in subtopics:
            for bloom_level in ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]:
                print("\n" + "=" * 80)
                print(f"🔹 Topic: {topic} | Subtopic: {subtopic} | Bloom: {bloom_level}")
                print("=" * 80)

                for q_num in range(1, QUESTIONS_PER_CELL + 1):
                    print(f"\n➡️ Generating question {q_num}...")
                    try:
                        row = run_pipeline(
                            openai_client,
                            anthropic_client,
                            slide_chunks,
                            slide_embeddings,
                            chunk_sources,
                            topic,
                            subtopic,
                            bloom_level,
                            SUBJECT,
                            OUTPUT_CSV,
                        )
                        if row:
                            print("✅ Generated question:")
                            print("   ", row["question_stem"])
                            print("   Correct:", row["correct_answer"])
                            total_questions += 1
                        else:
                            print("❌ Question generation failed.")
                    except Exception as e:
                        print("🔥 Critical error during generation:", e)
                    time.sleep(2)  # rate limiting buffer

    print("\n🎉 Done!")
    print("Total questions generated:", total_questions)
    print("CSV saved to:", OUTPUT_CSV)

if __name__ == "__main__":
    main()




























#!/usr/bin/env python3
"""
Slide-scoped Java MCQ generator for COMP1050.

- Reads your lecture PPTX files
- Chunks & embeds them with OpenAI embeddings
- Retrieves relevant slide chunks per concept
- Generates MCQs with OpenAI (gpt-4.1 by default)
- Evaluates each question with Claude 4 Haiku
- Writes everything (including retrieved chunks) to a CSV
- DEBUG mode prints the exact slide chunks used per question

Dependencies:
    pip install python-pptx openai anthropic numpy pandas
Environment:
    export OPENAI_API_KEY="..."
    export CLAUDE_ANTHROPIC_KEY="..."
"""

import os
import re
import csv
import uuid
import time
from datetime import datetime

import numpy as np
import pandas as pd
from pptx import Presentation
from openai import OpenAI
from anthropic import Anthropic

# ==========================
# CONFIG / CONSTANTS
# ==========================

# Toggle verbose printing of retrieved slide chunks
DEBUG = True

# OpenAI + Anthropic models
OPENAI_CHAT_MODEL = "gpt-4.1"
OPENAI_EMBED_MODEL = "text-embedding-3-small"
CLAUDE_MODEL = "claude-haiku-4-5-20251001"

# Where to save questions
OUTPUT_CSV = "java_questions_slide_scoped.csv"

# How many questions per (concept, bloom_level)
QUESTIONS_PER_CELL = 2

# Topics and Bloom levels
TOPICS = [
    "Java fundamentals",
    "Objects and Classes",
    "Inheritance",
    "Polymorphism",
    "Abstract Classes",
    "Interfaces",
    "Generics",
    "Java Collections Framework",
    "JavaFX basics",
    "Event handling in JavaFX",
]

BLOOM_LEVELS = ["Remember", "Understand", "Apply", "Analyze", "Evaluate"]
SUBJECT = "Java"

# Slide PPTX paths (relative to this script directory)
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
# SLIDE_PATHS = [
#     os.path.join(SCRIPT_DIR, "l1_review.pptx"),
#     os.path.join(SCRIPT_DIR, "l2_objects.pptx"),
#     os.path.join(SCRIPT_DIR, "l3_abstract.pptx"),
#     os.path.join(SCRIPT_DIR, "l4_interfaces.pptx"),
#     os.path.join(SCRIPT_DIR, "l6_-javaFX.pptx"),
#     os.path.join(SCRIPT_DIR, "l5_generics.pptx"),
#     os.path.join(SCRIPT_DIR, "comp1050-lecture9-JCF.pptx"),
#     os.path.join(SCRIPT_DIR, "l7_event.pptx"),
#     os.path.join(SCRIPT_DIR, "l2b_classes.pptx"),
# ]



SLIDES_FOLDER = os.path.join(SCRIPT_DIR, "slides")  # or any folder you want

def get_all_pptx_files(folder):
    ppt_files = []
    for fname in os.listdir(folder):
        if fname.lower().endswith(".pptx"):
            ppt_files.append(os.path.join(folder, fname))
    if not ppt_files:
        raise RuntimeError(f"No PPTX slide files found in {folder}")
    return ppt_files

SLIDE_PATHS = get_all_pptx_files(SLIDES_FOLDER)

print("📄 Found slide files:")
for p in SLIDE_PATHS:
    print("   -", p)

# ==========================
# API CLIENTS
# ==========================

def init_clients():
    openai_key = os.getenv("OPENAI_API_KEY")
    claude_key = os.getenv("CLAUDE_ANTHROPIC_KEY")

    if not openai_key:
        raise RuntimeError("OPENAI_API_KEY environment variable is not set.")
    if not claude_key:
        raise RuntimeError("CLAUDE_ANTHROPIC_KEY environment variable is not set.")

    openai_client = OpenAI(api_key=openai_key)
    anthropic_client = Anthropic(api_key=claude_key)
    return openai_client, anthropic_client


# ==========================
# SLIDES → TEXT
# ==========================

def ppt_to_text(path: str) -> str:
    """Extract visible text from a PPTX file."""
    if not os.path.exists(path):
        print(f"⚠️ Slide file not found: {path}")
        return ""

    prs = Presentation(path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def load_all_slides(paths) -> str:
    """Combine text from all slide decks into a single string."""
    all_texts = []
    for p in paths:
        print(f"📥 Loading slides from: {p}")
        t = ppt_to_text(p)
        if t.strip():
            all_texts.append(t)
    joined = "\n\n".join(all_texts)
    print("✅ Total slide text length:", len(joined))
    return joined


# ==========================
# CHUNKING & EMBEDDINGS
# ==========================

def chunk_text(text: str, max_len: int = 1200, overlap: int = 200):
    """Simple char-based chunking with overlap and soft breaks on line boundaries."""
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(n, start + max_len)
        chunk = text[start:end]
        # try to break at a newline if possible
        last_break = max(chunk.rfind("\n\n"), chunk.rfind("\n"))
        if last_break > 200:  # don't break too early
            end = start + last_break
            chunk = text[start:end]
        chunk = chunk.strip()
        if len(chunk) > 50:
            chunks.append(chunk)
        start = max(start + 1, end - overlap)
    print(f"📌 Number of chunks: {len(chunks)}")
    return chunks


def embed_texts(openai_client: OpenAI, texts):
    """Get embeddings for a list of texts using OpenAI."""
    if not texts:
        return []
    resp = openai_client.embeddings.create(
        model=OPENAI_EMBED_MODEL,
        input=texts
    )
    return [np.array(d.embedding, dtype=np.float32) for d in resp.data]


def cosine_similarity(a: np.ndarray, b: np.ndarray) -> float:
    denom = np.linalg.norm(a) * np.linalg.norm(b)
    if denom == 0:
        return 0.0
    return float(np.dot(a, b) / denom)


# ==========================
# RETRIEVAL
# ==========================

def get_relevant_slide_context(
    openai_client: OpenAI,
    query: str,
    slide_chunks,
    slide_embeddings,
    k: int = 4
):
    """
    Retrieve top-k most relevant slide chunks for the given query.
    Returns (combined_context_str, [list_of_chunks]).
    """
    if not slide_chunks or not slide_embeddings:
        return "", []

    query_emb = embed_texts(openai_client, [query])[0]
    sims = [cosine_similarity(query_emb, emb) for emb in slide_embeddings]
    top_indices = np.argsort(sims)[-k:][::-1]
    selected_chunks = [slide_chunks[i] for i in top_indices]
    combined = "\n\n".join(selected_chunks)
    return combined, selected_chunks


# ==========================
# PROMPT BUILDING
# ==========================

def build_prompt(
    concept: str,
    bloom_level: str,
    subject: str,
    context: str,
    chunks_for_debug=None
) -> str:
    """Build the slide-restricted MCQ generation prompt."""

    if DEBUG and chunks_for_debug:
        print("\n" + "=" * 80)
        print(f"🔍 DEBUG: Retrieved slide chunks for concept: {concept}")
        print("=" * 80)
        for idx, ch in enumerate(chunks_for_debug, start=1):
            print(f"\n--- Slide Chunk #{idx} ---\n{ch}\n")
        print("=" * 80 + "\n")

    return f"""
You are generating a Java multiple-choice question STRICTLY based on the following lecture content:

[LECTURE_SLIDES]
{context}
[/LECTURE_SLIDES]

Rules:
- Use ONLY content inside [LECTURE_SLIDES].
- The question MUST be answerable directly from this content.
- Do NOT invent new APIs, Java features, or examples not present in the slides.
- Concept: "{concept}"
- Bloom level: "{bloom_level}"

Produce ONE multiple-choice Java question (at most 50 words in the question stem)
in this EXACT format:

Question: [Your question stem here]

Options:
A) [Option A text]
B) [Option B text]
C) [Option C text]
D) [Option D text]

Correct Answer: [Letter of correct option only - A, B, C, or D]

Option A explanation: [Brief explanation of why option A is correct or incorrect]
Option B explanation: [Brief explanation of why option B is correct or incorrect]
Option C explanation: [Brief explanation of why option C is correct or incorrect]
Option D explanation: [Brief explanation of why option D is correct or incorrect]

Main Explanation: [Brief overall explanation of the answer]

Tags: {concept}, {bloom_level}
"""


# ==========================
# PARSING MODEL OUTPUT
# ==========================

def parse_question_response(response: str) -> dict:
    """Parse MCQ text from the model into structured fields."""
    patterns = {
        'question': r"[Qq]uestion:\s*(.+?)\s*[Oo]ptions?:",
        'option_a': r"A\)\s*(.+?)(?=\n[B-D]\)|\nCorrect|$)",
        'option_b': r"B\)\s*(.+?)(?=\n[C-D]\)|\nCorrect|$)",
        'option_c': r"C\)\s*(.+?)(?=\nD\)|\nCorrect|$)",
        'option_d': r"D\)\s*(.+?)(?=\nCorrect|$)",
        'correct_answer': r"Correct Answer:\s*([A-D])",
        'option_a_explanation': r"Option A explanation:\s*(.+?)(?=\nOption B|\nMain Explanation|$)",
        'option_b_explanation': r"Option B explanation:\s*(.+?)(?=\nOption C|\nMain Explanation|$)",
        'option_c_explanation': r"Option C explanation:\s*(.+?)(?=\nOption D|\nMain Explanation|$)",
        'option_d_explanation': r"Option D explanation:\s*(.+?)(?=\nMain Explanation|$)",
        'main_explanation': r"Main Explanation:\s*(.+?)(?=\nTags:|$)",
        'tags': r"Tags:\s*(.+)",
    }

    extracted = {}
    for key, pattern in patterns.items():
        m = re.search(pattern, response, re.DOTALL)
        extracted[key] = m.group(1).strip() if m else ""
    return extracted


# ==========================
# GENERATION & EVALUATION
# ==========================

def generate_with_openai(openai_client: OpenAI, prompt: str) -> str:
    """Generate MCQ text via OpenAI ChatCompletion."""
    resp = openai_client.chat.completions.create(
        model=OPENAI_CHAT_MODEL,
        messages=[
            {
                "role": "system",
                "content": (
                    "You are a Java exam question generator. "
                    "You must obey the user's format exactly and use only the provided lecture content."
                )
            },
            {"role": "user", "content": prompt}
        ],
        temperature=0.8,
        top_p=0.95,
        max_tokens=500
    )
    return resp.choices[0].message.content.strip()


def build_eval_prompt(
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    concept: str,
) -> str:
    return f"""
You are evaluating the quality of the following Java multiple-choice question.

Question: \"\"\"{question}\"\"\"
Options:
A) {option_a}
B) {option_b}
C) {option_c}
D) {option_d}

Correct Answer: {answer}

Option A explanation: {a_explanation}
Option B explanation: {b_explanation}
Option C explanation: {c_explanation}
Option D explanation: {d_explanation}
Overall Explanation: {main_explanation}

Rate 1–5 on:
1. Relevance — Does it directly assess the concept "{concept}"?
2. Accuracy — Is the correct answer and distractors technically accurate and reasonable?
3. Explainability — Are explanations clear, specific, and helpful?
4. Bloom Alignment — Does the cognitive demand match Bloom: "{bloom_level}"?

Return exactly:
relevance=X,bloom=Y,accuracy=Z,explainability=W,"Justification: [text]"
"""


def parse_eval_response(response: str):
    pattern = (
        r"relevance[=:]\s*(\d).*?"
        r"bloom[=:]\s*(\d).*?"
        r"accuracy[=:]\s*(\d).*?"
        r"explainability[=:]\s*(\d).*?"
        r"Justification:\s*\"?(.+?)\"?$"
    )
    m = re.search(pattern, response, re.DOTALL | re.IGNORECASE)
    if not m:
        print("⚠️ Could not parse evaluation response:\n", response)
        return None

    try:
        return {
            "eval_relevance": int(m.group(1)),
            "eval_bloom_alignment": int(m.group(2)),
            "eval_accuracy": int(m.group(3)),
            "eval_explainability": int(m.group(4)),
            "eval_justification": m.group(5).strip(),
        }
    except Exception as e:
        print("⚠️ Error parsing eval scores:", e)
        return None


def evaluate_with_claude(
    anthropic_client: Anthropic,
    question: str,
    option_a: str,
    option_b: str,
    option_c: str,
    option_d: str,
    answer: str,
    a_explanation: str,
    b_explanation: str,
    c_explanation: str,
    d_explanation: str,
    main_explanation: str,
    bloom_level: str,
    concept: str,
):
    prompt = build_eval_prompt(
        question,
        option_a,
        option_b,
        option_c,
        option_d,
        answer,
        a_explanation,
        b_explanation,
        c_explanation,
        d_explanation,
        main_explanation,
        bloom_level,
        concept,
    )

    resp = anthropic_client.messages.create(
        model=CLAUDE_MODEL,
        max_tokens=300,
        temperature=0.4,
        messages=[{"role": "user", "content": prompt}],
    )
    text = resp.content[0].text
    return parse_eval_response(text)


# ==========================
# PIPELINE
# ==========================

def run_pipeline(
    openai_client: OpenAI,
    anthropic_client: Anthropic,
    slide_chunks,
    slide_embeddings,
    concept: str,
    bloom_level: str,
    subject: str,
    csv_path: str
):
    # Retrieval
    context, chunks = get_relevant_slide_context(
        openai_client, concept, slide_chunks, slide_embeddings, k=4
    )
    if not context.strip():
        print("⚠️ No context retrieved for:", concept)
        return None

    # Build prompt (with debug printing)
    prompt = build_prompt(concept, bloom_level, subject, context, chunks_for_debug=chunks)

    # Generate question
    raw_response = generate_with_openai(openai_client, prompt)
    question_data = parse_question_response(raw_response)

    required_keys = ["question", "option_a", "option_b", "option_c", "option_d", "correct_answer"]
    if not all(question_data.get(k, "").strip() for k in required_keys):
        print("❌ Incomplete question data, skipping.")
        return None

    # Evaluate with Claude
    eval_scores = evaluate_with_claude(
        anthropic_client,
        question=question_data["question"],
        option_a=question_data["option_a"],
        option_b=question_data["option_b"],
        option_c=question_data["option_c"],
        option_d=question_data["option_d"],
        answer=question_data["correct_answer"],
        a_explanation=question_data.get("option_a_explanation", ""),
        b_explanation=question_data.get("option_b_explanation", ""),
        c_explanation=question_data.get("option_c_explanation", ""),
        d_explanation=question_data.get("option_d_explanation", ""),
        main_explanation=question_data.get("main_explanation", ""),
        bloom_level=bloom_level,
        concept=concept,
    )

    retrieved_chunk_text = "\n\n===== CHUNK BREAK =====\n\n".join(chunks)

    row = {
        "id": str(uuid.uuid4()),
        "timestamp": datetime.now().isoformat(),
        "concept": concept,
        "bloom_level": bloom_level,
        "question_stem": question_data["question"],
        "option_a": question_data["option_a"],
        "option_b": question_data["option_b"],
        "option_c": question_data["option_c"],
        "option_d": question_data["option_d"],
        "correct_answer": question_data["correct_answer"],
        "a_explanation": question_data.get("option_a_explanation", ""),
        "b_explanation": question_data.get("option_b_explanation", ""),
        "c_explanation": question_data.get("option_c_explanation", ""),
        "d_explanation": question_data.get("option_d_explanation", ""),
        "main_explanation": question_data.get("main_explanation", ""),
        "raw_model_response": raw_response,
        "retrieved_chunks": retrieved_chunk_text,
    }

    if eval_scores:
        row.update(eval_scores)
    else:
        row.update(
            {
                "eval_relevance": None,
                "eval_bloom_alignment": None,
                "eval_accuracy": None,
                "eval_explainability": None,
                "eval_justification": None,
            }
        )

    # Append to CSV
    file_exists = os.path.exists(csv_path)
    with open(csv_path, "a", newline="", encoding="utf-8") as f:
        writer = csv.DictWriter(f, fieldnames=row.keys())
        if not file_exists:
            writer.writeheader()
        writer.writerow(row)

    return row


# ==========================
# MAIN
# ==========================

def main():
    # Init clients
    openai_client, anthropic_client = init_clients()

    # Load slide text & build chunks/embeddings
    raw_text = load_all_slides(SLIDE_PATHS)
    slide_chunks = chunk_text(raw_text, max_len=1200, overlap=200)

    print("🔍 Embedding slide chunks (one-time)...")
    slide_embeddings = embed_texts(openai_client, slide_chunks)
    print("✅ Embeddings ready.")

    print(f"💾 Output CSV will be: {OUTPUT_CSV}")
    questions_generated = 0

    for concept in TOPICS:
        for bloom_level in BLOOM_LEVELS:
            print("\n" + "=" * 80)
            print(f"🔹 Concept: {concept} | Bloom: {bloom_level}")
            print("=" * 80)

            for q_num in range(1, QUESTIONS_PER_CELL + 1):
                print(f"\n➡️ Generating Q{q_num}...")
                try:
                    row = run_pipeline(
                        openai_client,
                        anthropic_client,
                        slide_chunks,
                        slide_embeddings,
                        concept,
                        bloom_level,
                        SUBJECT,
                        OUTPUT_CSV,
                    )
                    if row:
                        print("✅ Question generated:")
                        print("   ", row["question_stem"])
                        print("   Correct:", row["correct_answer"])
                        questions_generated += 1
                    else:
                        print("❌ Generation failed for this question.")
                except Exception as e:
                    print("🔥 Critical error:", e)

                time.sleep(2)  # be nice to the APIs

    print("\n🎉 Done! Total questions generated:", questions_generated)
    print("📄 Saved to:", OUTPUT_CSV)


if __name__ == "__main__":
    main()
